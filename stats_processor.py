import os
import re
import json
import time
import datetime
import requests
import gc
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.dml.color import RGBColor
import config
import activity_logger

# --- [1. 유틸리티 함수군] ---

def normalize_title(title):
    if not title: return ""
    return re.sub(r'[^a-zA-Z0-9]', '', title).lower()

def find_all_templates(root_dir):
    templates = []
    if not os.path.isdir(root_dir): return templates
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith('.pptx') and not filename.startswith('~$'):
                templates.append(os.path.join(dirpath, filename))
    return templates

def find_templates_for_customer(template_base_dir, customer_path):
    customer_template_dir = os.path.join(template_base_dir, customer_path)
    return find_all_templates(customer_template_dir)

def find_all_panels_recursively(panel_list):
    all_panels = []
    for panel in panel_list:
        all_panels.append(panel)
        if panel.get("type") == "row" and "panels" in panel:
            all_panels.extend(find_all_panels_recursively(panel["panels"]))
    return all_panels

def get_dashboard_definition(dashboard_uid, retries=2, delay=1):
    url = f"{config.GRAFANA_URL.rstrip('/')}/api/dashboards/uid/{dashboard_uid}"
    headers = {"Authorization": f"Bearer {config.API_KEY}"}
    for i in range(retries):
        try:
            response = requests.get(url, headers=headers, timeout=30, verify=config.VERIFY_SSL)
            response.raise_for_status()
            dashboard = response.json().get('dashboard', {})
            dashboard['all_panels'] = find_all_panels_recursively(dashboard.get('panels', []))
            return dashboard
        except:
            if i < retries - 1: time.sleep(delay)
    return None

def find_panel_by_title(all_panels, title_from_placeholder):
    normalized_target = normalize_title(title_from_placeholder)
    for panel in all_panels:
        if normalize_title(panel.get('title', '')) == normalized_target:
            return panel
    return None

def get_all_placeholders(prs):
    placeholders = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    matches = re.findall(r'(\{\{.*?\}\})', p.text)
                    for match in matches: placeholders.add(match)
    return list(placeholders)

def replace_text_in_presentation(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for p in shape.text_frame.paragraphs:
                if "{{" not in p.text: continue
                original_run = p.runs[0] if p.runs else None
                original_font = original_run.font if original_run else None
                temp_text = p.text
                text_changed = False
                for ph, val in replacements.items():
                    if ph in temp_text:
                        temp_text = temp_text.replace(ph, str(val))
                        text_changed = True
                if text_changed:
                    p.text = temp_text
                    if original_font:
                        for run in p.runs:
                            try:
                                run.font.name = original_font.name
                                if original_font.size: run.font.size = original_font.size
                                run.font.bold = original_font.bold
                                if original_font.color and original_font.color.rgb:
                                    run.font.color.rgb = RGBColor.from_string(str(original_font.color.rgb))
                            except: pass

def calculate_previous_month_dates():
    today = datetime.date.today()
    end_date = today.replace(day=1) - relativedelta(days=1)
    start_date = end_date.replace(day=1)
    y_str, m_str = end_date.strftime("%Y"), end_date.strftime("%m")
    target_date_str = f"{y_str}.{int(m_str)}월"
    return {
        "placeholders": {
            "{{START_DATE}}": start_date.strftime("%Y년 %m월 %d일"),
            "{{END_DATE}}": end_date.strftime("%Y년 %m월 %d일"),
            "{{YEAR}}": y_str, "{{MONTH}}": m_str,
            "{{DATE_RANGE}}": f"{start_date.strftime('%Y년 %m월 %d일')} ~ {end_date.strftime('%Y년 %m월 %d일')}",
            "{{DATE_RANGE_HYPHEN}}": f"{start_date.strftime('%Y-%m-%d')} ~ {end_date.strftime('%Y-%m-%d')}"
        },
        "target_date_str": target_date_str,
        "start_ts": int(datetime.datetime.combine(start_date, datetime.time.min).timestamp() * 1000),
        "end_ts": int(datetime.datetime.combine(end_date, datetime.time.max).timestamp() * 1000)
    }

# --- [2. Grafana 데이터 추출 핵심] ---

def extract_stats_from_response(response_json, query_letter):
    try:
        frames = response_json.get('results', {}).get(query_letter, {}).get('frames', [])
        if not frames: return None
        fields = frames[0].get('schema', {}).get('fields', [])
        for idx, field in enumerate(fields):
            if field.get('type') in ['number', 'float', 'int']:
                values = frames[0].get('data', {}).get('values', [])[idx]
                valid = [v for v in values if v is not None]
                if valid:
                    return {"max": max(valid), "mean": sum(valid) / len(valid)}
    except: pass
    return None

def get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts):
    if not panel: return None
    headers = {"Authorization": f"Bearer {config.API_KEY}", "Content-Type": "application/json"}
    all_queries = panel.get('targets', [])
    selected = [q.copy() for q in all_queries if q.get('refId') == query_letter]
    if not selected: return None
    panel_ds = panel.get('datasource')
    for q in selected:
        if q.get('datasource') is None: q['datasource'] = panel_ds
        q.pop('real_hosts', None)
        q.setdefault('maxDataPoints', 720)
        q.setdefault('intervalMs', 3600000)
    try:
        response = requests.post(f"{config.GRAFANA_URL.rstrip('/')}/api/ds/query",
                                 headers=headers, data=json.dumps({"queries": selected, "from": str(start_ts), "to": str(end_ts)}),
                                 timeout=30, verify=config.VERIFY_SSL)
        response.raise_for_status()
        return extract_stats_from_response(response.json(), query_letter)
    except: return None

# --- [3. 기능 1: 비스트리밍 버전] ---

def process_statistics(customer_name=None):
    results = {"success": True, "processed_files": [], "errors": [], "failed_placeholders": []}
    date_info = calculate_previous_month_dates()
    target_date_str = date_info["target_date_str"]
    start_ts, end_ts = date_info["start_ts"], date_info["end_ts"]

    templates = find_templates_for_customer(config.OUTPUT_DIR_WITH_IMAGES, customer_name) if customer_name \
                else find_all_templates(config.OUTPUT_DIR_WITH_IMAGES)

    dashboard_defs_cache = {}
    panel_data_cache = {}
    current_customer = None

    for template_path in templates:
        filename = os.path.basename(template_path)
        relative_path = os.path.relpath(os.path.dirname(template_path), config.OUTPUT_DIR_WITH_IMAGES)
        folder = relative_path.split(os.sep)[0] if relative_path != '.' else "root"

        if current_customer != folder:
            panel_data_cache = {}
            current_customer = folder

        try:
            prs = Presentation(template_path)
            replacements = date_info["placeholders"].copy()
            grafana_query_count = 0
            file_failed_count = 0

            db_def = dashboard_defs_cache.get(folder, None)
            if db_def is None and folder != "root":
                uid = config.get_dashboard_uid(folder)
                if uid:
                    db_def = get_dashboard_definition(uid)
                    dashboard_defs_cache[folder] = db_def if db_def else False
                else:
                    dashboard_defs_cache[folder] = False

            if db_def:
                all_ph = get_all_placeholders(prs)
                for ph in all_ph:
                    if ph in replacements: continue
                    parts = ph.strip("{}").rsplit('_', 1)
                    if len(parts) == 2:
                        cache_key = (parts[0], parts[1])
                        if cache_key not in panel_data_cache:
                            panel = find_panel_by_title(db_def.get('all_panels', []), parts[0].replace("-", " "))
                            panel_data_cache[cache_key] = get_grafana_stats_by_panel(panel, parts[1], start_ts, end_ts)

                        stats = panel_data_cache[cache_key]
                        if stats:
                            replacements[ph] = config.SENTENCE_TEMPLATE.format(max=f"{stats['max']:.1f}", mean=f"{stats['mean']:.1f}")
                            grafana_query_count += 1
                        else:
                            replacements[ph] = "N/A"
                            file_failed_count += 1
                            results["failed_placeholders"].append({"file": filename, "placeholder": ph})

            replace_text_in_presentation(prs, replacements)

            new_filename = re.sub(r'\d{4}[.년]\s*\d{1,2}월', target_date_str, filename).replace(" - 복사본", "")
            output_subdir = os.path.join(config.OUTPUT_DIR, relative_path)
            os.makedirs(output_subdir, exist_ok=True)
            prs.save(os.path.join(output_subdir, new_filename))

            results["processed_files"].append({
                "template": filename,
                "customer": folder,
                "new_filename": new_filename,
                "grafana_queries": grafana_query_count,
                "failed_placeholders": file_failed_count
            })

            del prs
            gc.collect()
        except Exception as e:
            results["errors"].append(f"{filename}: {str(e)}")
            results["success"] = False

    return results

# --- [4. 기능 2: 스트리밍 버전] ---

def process_statistics_stream(customer_name=None):
    log_id = activity_logger.create_log("통계 삽입", customer_name)
    processed_files = []
    errors = []
    failed_placeholders = []
    log_completed = False

    try:
        date_info = calculate_previous_month_dates()
        target_date_str = date_info["target_date_str"]
        start_ts, end_ts = date_info["start_ts"], date_info["end_ts"]
        templates = find_templates_for_customer(config.OUTPUT_DIR_WITH_IMAGES, customer_name) if customer_name \
                    else find_all_templates(config.OUTPUT_DIR_WITH_IMAGES)

        total = len(templates)
        yield {"type": "init", "total_files": total, "customer": customer_name or "전체"}

        dashboard_defs_cache = {}
        panel_data_cache = {}
        current_customer = None

        for idx, template_path in enumerate(templates, 1):
            filename = os.path.basename(template_path)
            relative_path = os.path.relpath(os.path.dirname(template_path), config.OUTPUT_DIR_WITH_IMAGES)
            folder = relative_path.split(os.sep)[0] if relative_path != '.' else "root"

            if current_customer != folder:
                panel_data_cache = {}
                current_customer = folder

            yield {"type": "file_start", "file_index": idx, "total_files": total, "filename": filename, "customer": folder}
            file_logs = [f"[{idx}/{total}] {folder}/{filename} 처리 시작"]

            try:
                prs = Presentation(template_path)
                final_replacements = date_info["placeholders"].copy()
                grafana_query_count = 0
                file_failed_count = 0

                db_def = dashboard_defs_cache.get(folder, None)
                skip_msg = None
                if db_def is None and folder != "root":
                    uid = config.get_dashboard_uid(folder)
                    if uid:
                        file_logs.append(f"  → 대시보드 로딩: {uid}")
                        db_def = get_dashboard_definition(uid)
                        if db_def:
                            dashboard_defs_cache[folder] = db_def
                            file_logs.append(f"  ✅ 대시보드 로딩 완료 (패널 {len(db_def.get('all_panels', []))}개)")
                        else:
                            dashboard_defs_cache[folder] = False
                            skip_msg = f"❌ 대시보드 로딩 실패 — 통계 삽입 건너뜀 ({uid})"
                            file_logs.append(f"  {skip_msg}")
                    else:
                        dashboard_defs_cache[folder] = False
                        skip_msg = f"⚠️ 대시보드 매핑 없음 — 통계 삽입 건너뜀 ('{folder}')"
                        file_logs.append(f"  {skip_msg}")
                elif db_def is False:
                    skip_msg = "⏭️ 대시보드 없음 (캐시) — 통계 삽입 건너뜀"
                    file_logs.append(f"  {skip_msg}")

                if not db_def and skip_msg:
                    yield {
                        "type": "file_progress",
                        "file_index": idx, "total_files": total,
                        "filename": filename, "customer": folder,
                        "placeholder_current": 0, "placeholder_total": 0,
                        "message": skip_msg
                    }

                if db_def:
                    all_ph = get_all_placeholders(prs)
                    grafana_phs = [ph for ph in all_ph if ph not in final_replacements
                                   and len(ph.strip("{}").rsplit('_', 1)) == 2]
                    ph_total = len(grafana_phs)
                    file_logs.append(f"  → 플레이스홀더 {len(all_ph)}개 발견 (Grafana 대상: {ph_total}개)")
                    ph_idx = 0

                    for ph in all_ph:
                        if ph in final_replacements: continue
                        parts = ph.strip("{}").rsplit('_', 1)
                        if len(parts) == 2:
                            ph_idx += 1
                            panel_title_raw = parts[0]
                            query_letter = parts[1]
                            cache_key = (panel_title_raw, query_letter)

                            if cache_key not in panel_data_cache:
                                panel = find_panel_by_title(db_def.get('all_panels', []), panel_title_raw.replace("-", " "))
                                if panel is None:
                                    msg = f"⚠️ 패널 없음: '{panel_title_raw}'"
                                    file_logs.append(f"  {msg}")
                                    panel_data_cache[cache_key] = None
                                    yield {
                                        "type": "file_progress",
                                        "file_index": idx, "total_files": total,
                                        "filename": filename, "customer": folder,
                                        "placeholder_current": ph_idx, "placeholder_total": ph_total,
                                        "message": msg
                                    }
                                else:
                                    panel_title_display = panel.get('title', panel_title_raw)
                                    msg = f"Grafana 조회: '{panel_title_display}' ({query_letter})"
                                    file_logs.append(f"  → {msg}")
                                    yield {
                                        "type": "file_progress",
                                        "file_index": idx, "total_files": total,
                                        "filename": filename, "customer": folder,
                                        "placeholder_current": ph_idx, "placeholder_total": ph_total,
                                        "message": msg
                                    }
                                    result = get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts)
                                    panel_data_cache[cache_key] = result
                                    if result:
                                        result_msg = f"✅ {panel_title_display} ({query_letter}): 최대 {result['max']:.1f} 평균 {result['mean']:.1f}"
                                    else:
                                        result_msg = f"❌ {ph}: 유효 데이터 없음"
                                    file_logs.append(f"  {result_msg}")
                                    yield {
                                        "type": "file_progress",
                                        "file_index": idx, "total_files": total,
                                        "filename": filename, "customer": folder,
                                        "placeholder_current": ph_idx, "placeholder_total": ph_total,
                                        "message": result_msg
                                    }
                            else:
                                cached = panel_data_cache[cache_key]
                                if cached:
                                    msg = f"♻️ 캐시 사용: '{panel_title_raw}' ({query_letter})"
                                    file_logs.append(f"  {msg}")
                                    yield {
                                        "type": "file_progress",
                                        "file_index": idx, "total_files": total,
                                        "filename": filename, "customer": folder,
                                        "placeholder_current": ph_idx, "placeholder_total": ph_total,
                                        "message": msg
                                    }

                            stats = panel_data_cache[cache_key]
                            if stats:
                                final_replacements[ph] = config.SENTENCE_TEMPLATE.format(max=f"{stats['max']:.1f}", mean=f"{stats['mean']:.1f}")
                                grafana_query_count += 1
                            else:
                                final_replacements[ph] = "N/A"
                                file_failed_count += 1
                                failed_placeholders.append({"file": filename, "placeholder": ph})

                replace_text_in_presentation(prs, final_replacements)
                new_filename = re.sub(r'\d{4}[.년]\s*\d{1,2}월', target_date_str, filename).replace(" - 복사본", "")
                output_path = os.path.join(config.OUTPUT_DIR, relative_path, new_filename)
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                prs.save(output_path)

                processed_files.append({
                    "template": filename,
                    "customer": folder,
                    "new_filename": new_filename,
                    "grafana_queries": grafana_query_count,
                    "failed_placeholders": file_failed_count
                })

                file_logs.append(f"  ✅ 완료: {grafana_query_count}개 삽입, 실패 {file_failed_count}개 → {new_filename}")
                activity_logger.add_detail_batch(log_id, file_logs)
                yield {
                    "type": "file_done", "file_index": idx, "total_files": total,
                    "success": True, "new_filename": new_filename,
                    "grafana_queries": grafana_query_count,
                    "failed_placeholders": file_failed_count
                }

                del prs
                gc.collect()
            except Exception as e:
                file_logs.append(f"  ❌ 오류: {str(e)}")
                activity_logger.add_detail_batch(log_id, file_logs)
                errors.append(f"{filename}: {str(e)}")
                yield {"type": "file_done", "file_index": idx, "total_files": total, "success": False, "error": str(e)}

        summary = {
            "processed_count": len(processed_files),
            "error_count": len(errors),
            "total_grafana_queries": sum(f.get("grafana_queries", 0) for f in processed_files)
        }
        activity_logger.complete_log(log_id, len(processed_files) > 0, summary)
        log_completed = True

        yield {
            "type": "complete",
            "success": len(processed_files) > 0,
            "processed_count": len(processed_files),
            "error_count": len(errors),
            "total_grafana_queries": summary["total_grafana_queries"],
            "failed_placeholder_count": len(failed_placeholders),
            "processed_files": processed_files,
            "errors": errors
        }
    except Exception as e:
        if not log_completed:
            activity_logger.complete_log(log_id, False, {"error": str(e)})
        yield {"type": "error", "error": str(e)}
