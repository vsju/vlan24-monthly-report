import os
import re
import datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import config
import activity_logger

def normalize_name(name):
    """이름에서 모든 특수문자와 공백을 제거하여 비교 가능한 형태로 만듭니다."""
    if not name:
        return ""
    return re.sub(r'[^a-zA-Z0-9]', '', name).lower()

def get_target_shapes_from_template(prs):
    """
    프레젠테이션의 모든 슬라이드를 돌면서
    (이름, 도형객체) 튜플 리스트로 수집.
    그룹 도형 내부도 재귀적으로 탐색.
    """
    target_shapes = []

    def collect_shapes(shapes):
        for shape in shapes:
            if shape.name:
                target_shapes.append((shape.name, shape))
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                collect_shapes(shape.shapes)

    for slide in prs.slides:
        collect_shapes(slide.shapes)

    return target_shapes

def find_all_templates(root_dir):
    templates = []
    if not os.path.isdir(root_dir):
        return templates
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith('.pptx') and not filename.startswith('~$'):
                full_path = os.path.join(dirpath, filename)
                templates.append(full_path)
    return templates

def find_templates_for_customer(template_base_dir, customer_path):
    customer_template_dir = os.path.join(template_base_dir, customer_path)
    return find_all_templates(customer_template_dir)

def calculate_previous_month_dates():
    today = datetime.date.today()
    end_date = today.replace(day=1) - relativedelta(days=1)
    start_date = end_date.replace(day=1)
    start_date_str_kr = start_date.strftime("%Y년 %m월 %d일")
    end_date_str_kr = end_date.strftime("%Y년 %m월 %d일")
    start_date_str_hyphen = start_date.strftime("%Y-%m-%d")
    end_date_str_hyphen = end_date.strftime("%Y-%m-%d")
    return {
        "placeholders": {
            "{{START_DATE}}": start_date_str_kr,
            "{{END_DATE}}": end_date_str_kr,
            "{{MONTH}}": end_date.strftime("%m"),
            "{{DATE_RANGE}}": f"{start_date_str_kr} ~ {end_date_str_kr}",
            "{{DATE_RANGE_HYPHEN}}": f"{start_date_str_hyphen} ~ {end_date_str_hyphen}"
        },
        "filename_date": end_date.strftime("%Y-%m")
    }

def find_available_images(root_dir):
    image_map = {}
    if not os.path.isdir(root_dir):
        return None
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                name_without_extension = os.path.splitext(filename)[0]
                image_map[name_without_extension] = os.path.join(dirpath, filename)
    return image_map

def replace_text_in_presentation(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    for placeholder, value in replacements.items():
                        if placeholder in r.text:
                            r.text = r.text.replace(placeholder, str(value))

def process_images(customer_name=None):
    results = {
        "success": True,
        "processed_files": [],
        "errors": [],
        "logs": [],
        "summary": {}
    }
    
    def log(message):
        results["logs"].append(message)
    
    try:
        date_info = calculate_previous_month_dates()
        date_replacements = date_info["placeholders"]
        log("날짜 정보 계산 완료")
        
        if customer_name:
            templates = find_templates_for_customer(config.BASE_TEMPLATE_DIR, customer_name)
            results["summary"]["mode"] = f"특정 고객사: {customer_name}"
            log(f"고객사 '{customer_name}' 템플릿 검색 중...")
        else:
            templates = find_all_templates(config.BASE_TEMPLATE_DIR)
            results["summary"]["mode"] = "전체 고객사"
            log("전체 고객사 템플릿 검색 중...")
        
        if not templates:
            results["success"] = False
            results["errors"].append("처리할 템플릿이 없습니다.")
            log("❌ 처리할 템플릿이 없습니다.")
            return results
        
        results["summary"]["total_templates"] = len(templates)
        log(f"총 {len(templates)}개 템플릿 발견")
        
        for idx, template_path in enumerate(templates, 1):
            filename = os.path.basename(template_path)
            log(f"[{idx}/{len(templates)}] {filename} 처리 중...")
            try:
                prs = Presentation(template_path)
                
                replace_text_in_presentation(prs, date_replacements)
                log(f"  → 날짜 플레이스홀더 치환 완료")
                
                relative_path = os.path.relpath(os.path.dirname(template_path), config.BASE_TEMPLATE_DIR)
                image_search_dir = os.path.join(config.BASE_IMAGE_DIR, relative_path)
                if not os.path.isdir(image_search_dir):
                    parent_dir = os.path.dirname(image_search_dir)
                    if os.path.isdir(parent_dir):
                        image_search_dir = parent_dir
                
                available_images = find_available_images(image_search_dir)
                if available_images is None:
                    available_images = {}
                log(f"  → {len(available_images)}개 이미지 발견")
                
                normalized_image_map = {
                    normalize_name(k): k for k in available_images.keys()
                }
                
                target_shapes = get_target_shapes_from_template(prs)
                insert_count = 0
                skipped_shapes = []
                
                for shape_name, shape_object in target_shapes:
                    normalized_shape_name = normalize_name(shape_name)
                    
                    if normalized_shape_name in normalized_image_map:
                        original_image_name = normalized_image_map[normalized_shape_name]
                        image_path = available_images[original_image_name]
                        
                        left, top = shape_object.left, shape_object.top
                        width, height = shape_object.width, shape_object.height
                        
                        slide = shape_object.part.slide
                        sp = shape_object._sp
                        sp.getparent().remove(sp)
                        
                        slide.shapes.add_picture(image_path, left, top, width, height)
                        insert_count += 1
                    else:
                        skipped_shapes.append(shape_name)
                
                output_subdir = os.path.join(config.OUTPUT_DIR_WITH_IMAGES, relative_path)
                if not os.path.exists(output_subdir):
                    os.makedirs(output_subdir)
                
                intermediate_output_path = os.path.join(
                    output_subdir,
                    os.path.basename(template_path)
                )
                prs.save(intermediate_output_path)
                log(f"  → {insert_count}개 이미지 삽입 완료")
                log(f"  ✅ 저장 완료: {os.path.basename(intermediate_output_path)}")
                
                results["processed_files"].append({
                    "template": os.path.basename(template_path),
                    "customer": relative_path if relative_path != '.' else 'root',
                    "images_inserted": insert_count,
                    "skipped_shapes": skipped_shapes,
                    "output_path": intermediate_output_path
                })
                
            except Exception as e:
                log(f"  ❌ 오류: {str(e)}")
                results["errors"].append(f"템플릿 처리 실패 ({os.path.basename(template_path)}): {str(e)}")
        
        results["summary"]["processed_count"] = len(results["processed_files"])
        results["summary"]["error_count"] = len(results["errors"])
        
        log(f"===== 작업 완료 =====")
        log(f"처리된 파일: {results['summary']['processed_count']}개")
        log(f"오류: {results['summary']['error_count']}개")
        
        if results["summary"]["processed_count"] == 0:
            results["success"] = False
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"전체 프로세스 실패: {str(e)}")
        log(f"❌ 전체 프로세스 실패: {str(e)}")
    
    return results


def process_images_stream(customer_name=None):
    log_id = activity_logger.create_log("이미지 삽입", customer_name)
    processed_files = []
    errors = []

    try:
        date_info = calculate_previous_month_dates()
        date_replacements = date_info["placeholders"]

        if customer_name:
            templates = find_templates_for_customer(config.BASE_TEMPLATE_DIR, customer_name)
        else:
            templates = find_all_templates(config.BASE_TEMPLATE_DIR)

        if not templates:
            msg = "처리할 템플릿이 없습니다."
            activity_logger.add_error(log_id, msg)
            activity_logger.complete_log(log_id, success=False, summary={"processed_count": 0, "error_count": 1})
            yield {"type": "error", "error": msg}
            return

        total = len(templates)
        yield {"type": "init", "total_files": total, "customer": customer_name or "전체"}
        activity_logger.add_detail(log_id, f"총 {total}개 템플릿 발견")

        for idx, template_path in enumerate(templates, 1):
            filename = os.path.basename(template_path)
            relative_path = os.path.relpath(os.path.dirname(template_path), config.BASE_TEMPLATE_DIR)
            file_customer = relative_path.split(os.sep)[0] if relative_path != '.' else 'root'

            yield {
                "type": "file_start",
                "file_index": idx,
                "total_files": total,
                "filename": filename,
                "customer": file_customer
            }
            activity_logger.add_detail(log_id, f"[{idx}/{total}] {file_customer}/{filename} 처리 시작")

            try:
                prs = Presentation(template_path)
                replace_text_in_presentation(prs, date_replacements)

                image_search_dir = os.path.join(config.BASE_IMAGE_DIR, relative_path)
                if not os.path.isdir(image_search_dir):
                    parent_dir = os.path.dirname(image_search_dir)
                    if os.path.isdir(parent_dir):
                        image_search_dir = parent_dir

                available_images = find_available_images(image_search_dir)
                if available_images is None:
                    available_images = {}

                normalized_image_map = {normalize_name(k): k for k in available_images.keys()}
                target_shapes = get_target_shapes_from_template(prs)
                insert_count = 0
                skipped_shapes = []
                total_shapes = len(target_shapes)

                for s_idx, (shape_name, shape_object) in enumerate(target_shapes, 1):
                    normalized_shape_name = normalize_name(shape_name)

                    if normalized_shape_name in normalized_image_map:
                        original_image_name = normalized_image_map[normalized_shape_name]
                        image_path = available_images[original_image_name]

                        left, top = shape_object.left, shape_object.top
                        width, height = shape_object.width, shape_object.height

                        slide = shape_object.part.slide
                        sp = shape_object._sp
                        sp.getparent().remove(sp)
                        slide.shapes.add_picture(image_path, left, top, width, height)
                        insert_count += 1

                        if s_idx % 5 == 0 or s_idx == total_shapes:
                            yield {
                                "type": "file_progress",
                                "file_index": idx,
                                "total_files": total,
                                "filename": filename,
                                "customer": file_customer,
                                "shape_current": s_idx,
                                "shape_total": total_shapes,
                                "shape_name": shape_name,
                                "inserted": insert_count
                            }
                    else:
                        skipped_shapes.append(shape_name)

                output_subdir = os.path.join(config.OUTPUT_DIR_WITH_IMAGES, relative_path)
                if not os.path.exists(output_subdir):
                    os.makedirs(output_subdir)

                intermediate_output_path = os.path.join(output_subdir, os.path.basename(template_path))
                prs.save(intermediate_output_path)

                processed_files.append({
                    "template": filename,
                    "customer": file_customer,
                    "images_inserted": insert_count,
                    "skipped_count": len(skipped_shapes),
                    "output_path": intermediate_output_path
                })

                yield {
                    "type": "file_done",
                    "file_index": idx,
                    "total_files": total,
                    "filename": filename,
                    "customer": file_customer,
                    "success": True,
                    "images_inserted": insert_count,
                    "skipped": len(skipped_shapes)
                }
                activity_logger.add_detail(log_id, f"  ✅ {filename}: {insert_count}개 이미지 삽입, {len(skipped_shapes)}개 스킵")

            except Exception as e:
                err_msg = f"템플릿 처리 실패 ({filename}): {str(e)}"
                errors.append(err_msg)
                activity_logger.add_error(log_id, err_msg)
                yield {
                    "type": "file_done",
                    "file_index": idx,
                    "total_files": total,
                    "filename": filename,
                    "customer": file_customer,
                    "success": False,
                    "error": str(e)
                }

        summary = {
            "processed_count": len(processed_files),
            "error_count": len(errors),
            "total_images_inserted": sum(f.get("images_inserted", 0) for f in processed_files)
        }
        success = len(processed_files) > 0
        activity_logger.complete_log(log_id, success=success, summary=summary)

        yield {
            "type": "complete",
            "success": success,
            "processed_count": summary["processed_count"],
            "error_count": summary["error_count"],
            "total_images_inserted": summary["total_images_inserted"],
            "processed_files": processed_files,
            "errors": errors,
            "log_id": log_id
        }

    except Exception as e:
        activity_logger.add_error(log_id, f"전체 프로세스 실패: {str(e)}")
        activity_logger.complete_log(log_id, success=False, summary={"processed_count": 0, "error_count": 1})
        yield {"type": "error", "error": f"전체 프로세스 실패: {str(e)}", "log_id": log_id}
