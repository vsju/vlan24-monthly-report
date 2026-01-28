"""
템플릿 자동 생성 모듈
- 마스터 템플릿 기반 고객사별 템플릿 생성
- VM 정보 기반 슬라이드 복제 및 플레이스홀더 치환
"""
import os
import re
import copy
import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

import config

PATTERN_SLIDE_FIRST_VM = 4
PATTERN_SLIDE_TWO_RESOURCES = 5
PATTERN_SLIDE_ONE_RESOURCE = 6
PATTERN_SLIDE_OTHER_VM_FIRST = 7
MIN_SLIDES_REQUIRED = 11

TRAILING_SLIDES_COUNT = 3

def validate_master_template(prs):
    """마스터 템플릿 유효성 검증"""
    if len(prs.slides) < MIN_SLIDES_REQUIRED:
        return False, f"마스터 템플릿에 최소 {MIN_SLIDES_REQUIRED}개 슬라이드가 필요합니다. (현재: {len(prs.slides)}개)"
    return True, None

def parse_vm_directory(dir_name):
    """디렉토리명에서 VM명과 IP 추출
    
    지원 형식:
    - VM명 (IP)  - 공백 구분
    - VM명_(IP)  - 언더스코어 구분
    - VM명(IP)   - 구분자 없음
    
    IP 형식: x.x.x.x (각 x는 1-3자리 숫자)
    """
    import re
    
    pattern = r'^(.+?)[ _]?\((\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3})\)$'
    match = re.match(pattern, dir_name)
    
    if match:
        return match.group(1), match.group(2)
    
    return dir_name, ""

def parse_image_filename(filename):
    """이미지 파일명에서 VM명, 리소스명, 패널ID 추출"""
    name = os.path.splitext(filename)[0]
    if '_:' not in name:
        return None, None, None
    
    parts = name.split('_:')
    vm_part = parts[0]
    resource_part = parts[1] if len(parts) > 1 else ""
    
    if '_' in resource_part:
        resource_name, panel_id = resource_part.rsplit('_', 1)
    else:
        resource_name = resource_part
        panel_id = ""
    
    resource_name = resource_name.lstrip('_')
    
    return vm_part, resource_name, panel_id

def sort_resources(resources):
    """리소스를 CPU → Memory → 디스크/디렉토리(이름순) 순서로 정렬"""
    def resource_sort_key(res):
        name = res['name'].upper()
        if name == 'CPU':
            return (0, '')
        elif name == 'MEMORY':
            return (1, '')
        else:
            return (2, res['name'])
    return sorted(resources, key=resource_sort_key)

def get_query_type(resource_name):
    """리소스명에 따른 쿼리 타입 결정"""
    if resource_name.upper() == "CPU":
        return "A"
    return "C"

def format_resource_label(resource_name):
    """리소스 레이블 포맷: CPU/Memory 외에는 뒤에 ':' 추가"""
    name_upper = resource_name.upper()
    if name_upper in ("CPU", "MEMORY"):
        return resource_name
    return f"{resource_name}:"

def analyze_customer_images(customer_name):
    """고객사 이미지 폴더 분석하여 VM 및 리소스 정보 추출"""
    customer_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
    
    if not os.path.exists(customer_dir):
        return None, f"고객사 폴더가 존재하지 않습니다: {customer_dir}"
    
    vms = []
    for item in sorted(os.listdir(customer_dir)):
        item_path = os.path.join(customer_dir, item)
        if not os.path.isdir(item_path) or item.startswith('.'):
            continue
        
        vm_name, vm_ip = parse_vm_directory(item)
        
        images = [f for f in os.listdir(item_path) 
                 if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
        
        resources = []
        for img in sorted(images):
            vm_part, resource_name, panel_id = parse_image_filename(img)
            if resource_name:
                resources.append({
                    "name": resource_name,
                    "filename": os.path.splitext(img)[0],
                    "panel_id": panel_id,
                    "query": get_query_type(resource_name),
                    "image_vm_name": vm_part
                })
        
        resources = sort_resources(resources)
        
        vms.append({
            "dir_name": item,
            "vm_name": vm_name,
            "ip": vm_ip,
            "resources": resources,
            "pages_needed": math.ceil(len(resources) / 2) if resources else 0
        })
    
    return vms, None

def duplicate_slide_xml(prs, slide_index):
    """슬라이드 복제 (XML 기반) - 새 슬라이드를 마지막에 추가"""
    from copy import deepcopy
    
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    for shape in source_slide.shapes:
        el = shape._element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    
    return new_slide

def replace_text_in_runs(shape, replacements):
    """도형 내 모든 run에서 플레이스홀더 치환 (서식 최대한 보존)
    
    전략:
    1. 각 run 내에서 완전한 플레이스홀더 치환 (서식 100% 보존)
    2. 분리된 플레이스홀더는 run 경계를 추적하여 치환 텍스트를 첫 run에 배치
       (분리된 플레이스홀더 영역만 병합, 나머지 run은 유지)
    """
    if not shape.has_text_frame:
        return False
    
    changed = False
    for para in shape.text_frame.paragraphs:
        runs = list(para.runs)
        if not runs:
            continue
        
        for run in runs:
            original = run.text
            new_text = original
            for old, new in replacements.items():
                if old in new_text:
                    new_text = new_text.replace(old, str(new))
            if new_text != original:
                run.text = new_text
                changed = True
        
        full_text = ''.join(run.text for run in runs)
        remaining_placeholders = []
        for placeholder in replacements.keys():
            if placeholder in full_text:
                remaining_placeholders.append(placeholder)
        
        if not remaining_placeholders:
            continue
        
        run_boundaries = []
        pos = 0
        for run in runs:
            run_boundaries.append((pos, pos + len(run.text), run))
            pos += len(run.text)
        
        for placeholder in remaining_placeholders:
            idx = full_text.find(placeholder)
            if idx == -1:
                continue
            
            start_run_idx = None
            end_run_idx = None
            for i, (start, end, run) in enumerate(run_boundaries):
                if start_run_idx is None and start <= idx < end:
                    start_run_idx = i
                if start < idx + len(placeholder) <= end:
                    end_run_idx = i
                    break
                if start_run_idx is not None and idx + len(placeholder) > end:
                    end_run_idx = i
            
            if start_run_idx is None or end_run_idx is None:
                continue
            
            if start_run_idx == end_run_idx:
                continue
            
            merged_text = ''.join(runs[i].text for i in range(start_run_idx, end_run_idx + 1))
            replacement_value = str(replacements[placeholder])
            new_merged = merged_text.replace(placeholder, replacement_value)
            
            runs[start_run_idx].text = new_merged
            for i in range(start_run_idx + 1, end_run_idx + 1):
                runs[i].text = ""
            
            changed = True
            
            full_text = ''.join(run.text for run in runs)
            pos = 0
            run_boundaries = []
            for run in runs:
                run_boundaries.append((pos, pos + len(run.text), run))
                pos += len(run.text)
    
    return changed


def set_text_preserve_format(shape, new_text):
    """도형 텍스트 전체를 새 텍스트로 교체 (서식 보존) - 단순 치환용"""
    if not shape.has_text_frame:
        return
    
    if not shape.text_frame.paragraphs:
        return
    
    first_para = shape.text_frame.paragraphs[0]
    if not first_para.runs:
        first_para.text = new_text
        return
    
    first_run = first_para.runs[0]
    first_run.text = new_text
    for run in list(first_para.runs)[1:]:
        run.text = ""
    
    for para in list(shape.text_frame.paragraphs)[1:]:
        for run in para.runs:
            run.text = ""

def replace_placeholders_in_text(text, replacements):
    """텍스트 내 플레이스홀더 치환"""
    result = text
    for old, new in replacements.items():
        result = result.replace(old, str(new))
    return result

def has_placeholder_in_shape(shape, placeholder):
    """도형 내 플레이스홀더 존재 여부 확인 (run 분리 케이스 포함)"""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        full_text = ''.join(run.text for run in para.runs)
        if placeholder in full_text:
            return True
    return False


def update_slide_content(slide, seq, vm_name, vm_ip, resources, log_func=None):
    """슬라이드 내용 업데이트 - 리소스별 플레이스홀더 치환
    
    처리 순서:
    1. 모든 도형에서 {{SEQ}}, {{VM}}, {{IP}} 치환
    2. 상단/하단 영역으로 나눠 {{RESOURCE}} 치환 (슬라이드 중간 12cm 기준)
    3. 이미지 도형(AUTO_SHAPE) 이름 변경
    
    영역 구분:
    - top < 12cm: 첫 번째 리소스 영역
    - top >= 12cm: 두 번째 리소스 영역
    """
    
    base_replacements = {
        '{{SEQ}}': seq,
        '{{IP}}': vm_ip,
    }
    
    BOUNDARY_CM = 12.0
    
    shapes_by_position = []
    for shape in slide.shapes:
        top_cm = shape.top / 914400 * 2.54
        shapes_by_position.append((top_cm, shape))
    shapes_by_position.sort(key=lambda x: x[0])
    
    for top_cm, shape in shapes_by_position:
        if shape.has_text_frame:
            replace_text_in_runs(shape, base_replacements)
    
    upper_resource_shapes = []
    lower_resource_shapes = []
    upper_image_shapes = []
    lower_image_shapes = []
    
    for top_cm, shape in shapes_by_position:
        if has_placeholder_in_shape(shape, '{{RESOURCE}}'):
            if top_cm < BOUNDARY_CM:
                upper_resource_shapes.append((top_cm, shape))
            else:
                lower_resource_shapes.append((top_cm, shape))
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if '{{RESOURCE_IMAGE}}' in shape.name:
                if top_cm < BOUNDARY_CM:
                    upper_image_shapes.append((top_cm, shape))
                else:
                    lower_image_shapes.append((top_cm, shape))
    
    if len(resources) >= 1:
        res = resources[0]
        image_vm = res.get('image_vm_name') or vm_name
        replacements = {
            '{{RESOURCE}}': format_resource_label(res['name']),
            '{{QUERY}}': res['query'],
            '{{VM}}': image_vm,
        }
        for top_cm, shape in upper_resource_shapes:
            replace_text_in_runs(shape, replacements)
        for top_cm, shape in upper_image_shapes:
            shape.name = res['filename']
            if log_func:
                log_func(f"    도형명 변경: {shape.name}")
    
    if len(resources) >= 2:
        res = resources[1]
        image_vm = res.get('image_vm_name') or vm_name
        replacements = {
            '{{RESOURCE}}': format_resource_label(res['name']),
            '{{QUERY}}': res['query'],
            '{{VM}}': image_vm,
        }
        for top_cm, shape in lower_resource_shapes:
            replace_text_in_runs(shape, replacements)
        for top_cm, shape in lower_image_shapes:
            shape.name = res['filename']
            if log_func:
                log_func(f"    도형명 변경: {shape.name}")

def set_cell_text_preserve_format(cell, new_text):
    """표 셀 텍스트를 변경하면서 서식 보존"""
    if not cell.text_frame.paragraphs:
        cell.text = new_text
        return
    
    first_para = cell.text_frame.paragraphs[0]
    if not first_para.runs:
        first_para.text = new_text
        return
    
    first_para.runs[0].text = new_text
    for run in list(first_para.runs)[1:]:
        run.text = ""


def update_table_vm_list(slide, vms):
    """슬라이드 4의 서버현황 표에 VM 목록 추가 (서식 보존, 행 자동 추가)
    
    - VM 개수에 맞게 테이블 행 자동 추가
    - 테이블 구조: 번호, VM명, IP, 비고 (4열)
    """
    from copy import deepcopy
    from pptx.table import _Row
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            tbl = table._tbl
            
            if len(tbl.tr_lst) < 2:
                continue
            
            template_row_idx = 1
            
            while len(tbl.tr_lst) > 2:
                tr = tbl.tr_lst[-1]
                tbl.remove(tr)
            
            for vm_idx, vm in enumerate(vms):
                if vm_idx == 0:
                    tr = tbl.tr_lst[template_row_idx]
                else:
                    template_tr = tbl.tr_lst[template_row_idx]
                    new_tr = deepcopy(template_tr)
                    tbl.append(new_tr)
                    tr = tbl.tr_lst[-1]
                
                row = _Row(tr, table)
                
                vm_name = vm.get('vm_name', '')
                vm_ip = vm.get('ip', '')
                seq_num = str(vm_idx + 1)
                
                cells = list(row.cells)
                if len(cells) > 0:
                    set_cell_text_preserve_format(cells[0], seq_num)
                if len(cells) > 1:
                    set_cell_text_preserve_format(cells[1], vm_name)
                if len(cells) > 2:
                    set_cell_text_preserve_format(cells[2], vm_ip)
            
            break

def generate_customer_template(master_template_path, customer_name, output_path=None, display_name=None):
    """고객사별 템플릿 자동 생성
    
    Args:
        master_template_path: 마스터 템플릿 경로
        customer_name: 고객사 폴더명 (name 필드)
        output_path: 출력 경로 (선택)
        display_name: 보고서에 표시할 이름 (없으면 customer_name 사용)
    """
    results = {
        "success": True,
        "logs": [],
        "errors": []
    }
    
    report_name = display_name if display_name else customer_name
    
    def log(msg):
        results["logs"].append(msg)
    
    try:
        vms, error = analyze_customer_images(customer_name)
        if error:
            results["success"] = False
            results["errors"].append(error)
            return results
        
        if not vms:
            results["success"] = False
            results["errors"].append("VM 폴더가 없습니다.")
            return results
        
        log(f"VM {len(vms)}개 발견")
        for vm in vms:
            log(f"  - {vm['vm_name']} ({vm['ip']}): 리소스 {len(vm['resources'])}개")
        
        prs = Presentation(master_template_path)
        total_slides = len(prs.slides)
        log(f"마스터 템플릿 로드: {total_slides}개 슬라이드")
        
        is_valid, validation_error = validate_master_template(prs)
        if not is_valid:
            results["success"] = False
            results["errors"].append(validation_error)
            return results
        
        for slide_idx in range(min(4, len(prs.slides))):
            slide = prs.slides[slide_idx]
            for shape in slide.shapes:
                if shape.has_text_frame and '{{CUSTOMER_NAME}}' in shape.text:
                    old_text = shape.text
                    new_text = old_text.replace('{{CUSTOMER_NAME}}', report_name)
                    set_text_preserve_format(shape, new_text)
        log(f"{{{{CUSTOMER_NAME}}}} 치환 완료: {report_name}")
        
        if len(prs.slides) > 3:
            update_table_vm_list(prs.slides[3], vms)
            log("슬라이드 4 서버현황표 업데이트 완료")
        
        for vm_idx, vm in enumerate(vms):
            seq = f"3.{vm_idx + 1}"
            vm_name = vm['vm_name']
            vm_ip = vm['ip']
            resources = vm['resources']
            
            log(f"VM {seq} {vm_name} 처리 중...")
            
            if not resources:
                log(f"  리소스 없음, 건너뜀")
                continue
            
            resource_pairs = []
            for i in range(0, len(resources), 2):
                pair = resources[i:i+2]
                resource_pairs.append(pair)
            
            for pair_idx, pair in enumerate(resource_pairs):
                if vm_idx == 0 and pair_idx == 0:
                    slide_template_idx = PATTERN_SLIDE_FIRST_VM
                elif pair_idx == 0:
                    slide_template_idx = PATTERN_SLIDE_OTHER_VM_FIRST
                elif len(pair) == 2:
                    slide_template_idx = PATTERN_SLIDE_TWO_RESOURCES
                else:
                    slide_template_idx = PATTERN_SLIDE_ONE_RESOURCE
                
                new_slide = duplicate_slide_xml(prs, slide_template_idx)
                
                update_slide_content(new_slide, seq, vm_name, vm_ip, pair, log)
                
                log(f"  슬라이드 생성: {', '.join([r['name'] for r in pair])}")
        
        trailing_slide_indices = list(range(
            PATTERN_SLIDE_OTHER_VM_FIRST + 1, 
            PATTERN_SLIDE_OTHER_VM_FIRST + 1 + TRAILING_SLIDES_COUNT
        ))
        
        trailing_slide_ids = []
        for idx in trailing_slide_indices:
            if idx < len(prs.slides):
                trailing_slide_ids.append(prs.slides._sldIdLst[idx])
        
        for sld_id in trailing_slide_ids:
            prs.slides._sldIdLst.remove(sld_id)
        
        for sld_id in trailing_slide_ids:
            prs.slides._sldIdLst.append(sld_id)
        
        log(f"마지막 {TRAILING_SLIDES_COUNT}개 슬라이드를 맨 뒤로 이동 완료")
        
        pattern_indices = sorted([PATTERN_SLIDE_FIRST_VM, PATTERN_SLIDE_TWO_RESOURCES, 
                                 PATTERN_SLIDE_ONE_RESOURCE, PATTERN_SLIDE_OTHER_VM_FIRST], 
                                reverse=True)
        
        for idx in pattern_indices:
            if idx < len(prs.slides):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
        
        log(f"패턴 슬라이드 4개 삭제 완료")
        
        if not output_path:
            output_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer_name)
            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, f"{customer_name}_월간보고서.pptx")
        
        prs.save(output_path)
        log(f"템플릿 저장: {output_path}")
        
        results["output_path"] = output_path
        results["vm_count"] = len(vms)
        results["slide_count"] = len(prs.slides)
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(str(e))
        import traceback
        results["traceback"] = traceback.format_exc()
    
    return results


def count_existing_vms_in_template(prs, log_func=None):
    """기존 템플릿에서 VM 개수 추정 (슬라이드 제목의 SEQ 번호 기반)
    
    슬라이드 제목 또는 첫 번째 텍스트 박스에서 '3.X' 형태의 SEQ 번호 찾기
    VM 슬라이드는 보통 "3.1 VM명" 또는 "3.12VM명" 형태의 제목을 가짐
    우선순위: 1) 제목 placeholder, 2) 일반 텍스트 도형 (줄 시작 매칭)
    """
    from pptx.enum.shapes import PP_PLACEHOLDER
    
    def log(msg):
        if log_func:
            log_func(msg)
    
    max_seq = 0
    title_pattern = re.compile(r'^\s*3\.(\d+)\b', re.MULTILINE)
    fallback_pattern = re.compile(r'^3\.(\d+)\b', re.MULTILINE)
    
    def get_shape_text(shape):
        """도형에서 전체 텍스트 추출"""
        if shape.has_text_frame:
            return shape.text.strip()
        return ""
    
    for slide_idx, slide in enumerate(prs.slides):
        title_text = None
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                        title_text = get_shape_text(shape)
                        log(f"슬라이드 {slide_idx + 1} 제목: '{title_text[:50] if title_text else '(없음)'}'")
                        break
                except:
                    pass
        
        if title_text:
            match = title_pattern.search(title_text)
            if match:
                seq_num = int(match.group(1))
                if seq_num > max_seq:
                    max_seq = seq_num
                    log(f"  → SEQ 3.{seq_num} 감지 (제목 placeholder)")
                continue
        
        for shape in slide.shapes:
            text = get_shape_text(shape)
            if text:
                match = fallback_pattern.search(text)
                if match:
                    seq_num = int(match.group(1))
                    if seq_num > max_seq:
                        max_seq = seq_num
                        log(f"슬라이드 {slide_idx + 1}: SEQ 3.{seq_num} 감지 (텍스트: {text[:40]})")
                    break
    
    log(f"최대 SEQ: 3.{max_seq} (VM {max_seq}개)")
    return max_seq


def insert_slide_at_position(prs, new_slide, position):
    """슬라이드를 특정 위치에 삽입 (0-indexed)
    
    새로 추가된 슬라이드(맨 뒤)를 원하는 위치로 이동
    """
    slide_id = prs.slides._sldIdLst[-1]
    prs.slides._sldIdLst.remove(slide_id)
    prs.slides._sldIdLst.insert(position, slide_id)


def add_vm_to_template(template_path, master_template_path, vm_dir_name, customer_name, seq_number=None):
    """기존 템플릿에 새 VM 슬라이드 추가
    
    Args:
        template_path: 기존 고객사 템플릿 경로
        master_template_path: 마스터 템플릿 경로 (패턴 슬라이드 원본)
        vm_dir_name: 추가할 VM 폴더명 (예: "VM명_(IP)")
        customer_name: 고객사 폴더명
        seq_number: SEQ 번호 (None이면 자동 계산)
    
    Returns:
        결과 dict (success, logs, errors, output_path 등)
    """
    results = {
        "success": True,
        "logs": [],
        "errors": []
    }
    
    def log(msg):
        results["logs"].append(msg)
    
    try:
        if not os.path.exists(template_path):
            results["success"] = False
            results["errors"].append(f"템플릿이 존재하지 않습니다: {template_path}")
            return results
        
        if not os.path.exists(master_template_path):
            results["success"] = False
            results["errors"].append(f"마스터 템플릿이 존재하지 않습니다: {master_template_path}")
            return results
        
        vm_path = os.path.join(config.BASE_IMAGE_DIR, customer_name, vm_dir_name)
        if not os.path.exists(vm_path):
            results["success"] = False
            results["errors"].append(f"VM 폴더가 존재하지 않습니다: {vm_path}")
            return results
        
        vm_name, vm_ip = parse_vm_directory(vm_dir_name)
        log(f"VM 정보: {vm_name} ({vm_ip})")
        
        images = [f for f in os.listdir(vm_path) 
                 if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
        
        resources = []
        for img in sorted(images):
            vm_part, resource_name, panel_id = parse_image_filename(img)
            if resource_name:
                resources.append({
                    "name": resource_name,
                    "filename": os.path.splitext(img)[0],
                    "panel_id": panel_id,
                    "query": get_query_type(resource_name),
                    "image_vm_name": vm_part
                })
        
        resources = sort_resources(resources)
        log(f"리소스 {len(resources)}개 발견: {', '.join([r['name'] for r in resources])}")
        
        if not resources:
            results["success"] = False
            results["errors"].append("리소스 이미지가 없습니다.")
            return results
        
        prs = Presentation(template_path)
        total_slides = len(prs.slides)
        log(f"기존 템플릿 로드: {total_slides}개 슬라이드")
        
        master_prs = Presentation(master_template_path)
        log(f"마스터 템플릿 로드: {len(master_prs.slides)}개 슬라이드")
        
        is_valid, validation_error = validate_master_template(master_prs)
        if not is_valid:
            results["success"] = False
            results["errors"].append(f"마스터 템플릿 오류: {validation_error}")
            return results
        
        if seq_number is None:
            existing_vm_count = count_existing_vms_in_template(prs, log_func=log)
            seq_number = existing_vm_count + 1
            log(f"기존 VM {existing_vm_count}개 감지, 새 SEQ: 3.{seq_number}")
        
        seq = f"3.{seq_number}"
        
        insert_position = total_slides - TRAILING_SLIDES_COUNT
        if insert_position < 0:
            insert_position = total_slides
        log(f"삽입 위치: 슬라이드 {insert_position + 1} (후행 {TRAILING_SLIDES_COUNT}개 슬라이드 앞)")
        
        resource_pairs = []
        for i in range(0, len(resources), 2):
            pair = resources[i:i+2]
            resource_pairs.append(pair)
        
        slides_added = 0
        for pair_idx, pair in enumerate(resource_pairs):
            if pair_idx == 0:
                pattern_idx = PATTERN_SLIDE_OTHER_VM_FIRST
            elif len(pair) == 2:
                pattern_idx = PATTERN_SLIDE_TWO_RESOURCES
            else:
                pattern_idx = PATTERN_SLIDE_ONE_RESOURCE
            
            source_slide = master_prs.slides[pattern_idx]
            slide_layout = source_slide.slide_layout
            
            matching_layout = None
            for layout in prs.slide_layouts:
                if layout.name == slide_layout.name:
                    matching_layout = layout
                    break
            
            if matching_layout is None:
                matching_layout = prs.slide_layouts[0]
            
            new_slide = prs.slides.add_slide(matching_layout)
            
            for shape in list(new_slide.shapes):
                sp = shape._element
                sp.getparent().remove(sp)
            
            from copy import deepcopy
            for shape in source_slide.shapes:
                el = shape._element
                new_el = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
            
            update_slide_content(new_slide, seq, vm_name, vm_ip, pair, log)
            
            current_position = len(prs.slides) - 1
            target_position = insert_position + slides_added
            
            slide_id = prs.slides._sldIdLst[current_position]
            prs.slides._sldIdLst.remove(slide_id)
            prs.slides._sldIdLst.insert(target_position, slide_id)
            
            slides_added += 1
            log(f"슬라이드 추가: {', '.join([r['name'] for r in pair])}")
        
        prs.save(template_path)
        log(f"템플릿 저장: {template_path}")
        
        results["output_path"] = template_path
        results["slides_added"] = slides_added
        results["vm_name"] = vm_name
        results["vm_ip"] = vm_ip
        results["seq"] = seq
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(str(e))
        import traceback
        results["traceback"] = traceback.format_exc()
    
    return results
