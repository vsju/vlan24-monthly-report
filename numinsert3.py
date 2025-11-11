from pptx import Presentation
from pptx.dml.color import RGBColor
import os
import datetime
from dateutil.relativedelta import relativedelta
import re
import requests
import json
import sys
import time
import gc

# ==============================================================================
# 1. 기본 설정 (보내주신 정보로 고정)
# ==============================================================================
# 이미지가 삽입된 파일들이 있는 폴더
BASE_TEMPLATE_DIR = "/root/Report/completed_with_images" 
# 최종 완성본을 저장할 폴더
OUTPUT_DIR = "/root/Report/completed_final"

# --- Grafana 연동 정보 ---
GRAFANA_URL = "http://localhost:3000"
API_KEY = "glsa_8RKuPs2USwIwxdafke3r4bcs93zkGO4E_462d232d"

# --- 고객사 폴더 이름과 대시보드 UID를 연결하는 지도 ---
DASHBOARD_MAP = {
    "kpmo": "dejkgjz0jnoqoa",
    "GIT": "aejkgkoze5nggb",
    "hansystem": "cejnb5yyuk5q8e",
    "humecca": "bejnb5db19blse",
    "klcns": "eejnb31cylreod",
    "sungwoo": "cejnb4aafury8e",
    "thepnl": "fejkgid897xtsc",
    "프리스타일": "fejkgfwux1fy8c"
}

# --- 자동 생성될 문장 형식 (수정 가능) ---
SENTENCE_TEMPLATE = "사용량 최대 {max}%, 평균 {mean}% 입니다."

# ==============================================================================
# 2. 자동화 코드 본문 (수정 금지)
# ==============================================================================
def normalize_title(title):
    if not title: return ""
    return re.sub(r'[^a-zA-Z0-9]', '', title).lower()

def find_all_templates(template_base_dir):
    template_paths = []
    if not os.path.isdir(template_base_dir): return template_paths
    for dirpath, _, filenames in os.walk(template_base_dir):
        for filename in filenames:
            if filename.lower().endswith('.pptx') and not filename.startswith('~$'):
                template_paths.append(os.path.join(dirpath, filename))
    return template_paths
    
def find_templates_for_customer(template_base_dir, customer_path):
    customer_template_dir = os.path.join(template_base_dir, customer_path)
    return find_all_templates(customer_template_dir)

def calculate_previous_month_dates():
    today = datetime.date.today()
    end_date = today.replace(day=1) - relativedelta(days=1)
    start_date = end_date.replace(day=1)
    start_date_str_kr = start_date.strftime("%Y년 %m월 %d일")
    end_date_str_kr = end_date.strftime("%Y년 %m월 %d일")
    start_date_str_hyphen = start_date.strftime("%Y-%m-%d")
    end_date_str_hyphen = end_date.strftime("%Y-%m-%d")
    return {
        "placeholders": { "{{START_DATE}}": start_date_str_kr, "{{END_DATE}}": end_date_str_kr, "{{MONTH}}": end_date.strftime("%m"), "{{DATE_RANGE}}": f"{start_date_str_kr} ~ {end_date_str_kr}", "{{DATE_RANGE_HYPHEN}}": f"{start_date_str_hyphen} ~ {end_date_str_hyphen}" },
        "filename_date": end_date.strftime("%Y-%m"),
        "start_ts": int(datetime.datetime.combine(start_date, datetime.time.min).timestamp() * 1000),
        "end_ts": int(datetime.datetime.combine(end_date, datetime.time.max).timestamp() * 1000)
    }

def find_all_panels_recursively(panel_list):
    all_panels = []
    for panel in panel_list:
        all_panels.append(panel)
        if panel.get("type") == "row" and "panels" in panel:
            all_panels.extend(find_all_panels_recursively(panel["panels"]))
    return all_panels

def get_dashboard_definition(dashboard_uid, retries=3, delay=1):
    url = f"{GRAFANA_URL}/api/dashboards/uid/{dashboard_uid}"
    headers = {"Authorization": f"Bearer {API_KEY}"}
    for i in range(retries):
        try:
            response = requests.get(url, headers=headers, timeout=20, verify=False)
            response.raise_for_status()
            dashboard = response.json().get('dashboard', {})
            dashboard['all_panels'] = find_all_panels_recursively(dashboard.get('panels', []))
            return dashboard
        except requests.exceptions.RequestException as e:
            print(f"  > 에러 (시도 {i+1}/{retries}): Grafana 대시보드({dashboard_uid}) 정보 가져오기 실패. {e}")
            time.sleep(delay)
    return None

def find_panel_by_title(all_panels, title_from_placeholder):
    normalized_placeholder_title = normalize_title(title_from_placeholder)
    for panel in all_panels:
        panel_title_from_grafana = panel.get('title', '')
        normalized_grafana_title = normalize_title(panel_title_from_grafana)
        if normalized_grafana_title and normalized_grafana_title == normalized_placeholder_title:
            return panel
    return None

def get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts):
    headers = {"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"}

    all_queries = panel.get('targets', [])
    # refId 일치하는 쿼리만 선택
    selected = [q.copy() for q in all_queries if q.get('refId') == query_letter]
    if not selected:
        return None

    panel_datasource = panel.get('datasource')
    if panel_datasource:
        for q in selected:
            if 'datasource' not in q or q['datasource'] is None:
                q['datasource'] = panel_datasource

    # 🔴 여기 한 줄 추가: real_hosts 제거
    for q in selected:
        q.pop('real_hosts', None)

    # 한 달치라도 포인트 수 제한 (예: 720개, 1시간 간격)
    for q in selected:
        q.setdefault('maxDataPoints', 720)
        q.setdefault('intervalMs', 3600000)

    query_payload = {
        "queries": selected,
        "from": str(start_ts),
        "to": str(end_ts)
    }
    query_url = f"{GRAFANA_URL}/api/ds/query"

    try:
        response = requests.post(query_url, headers=headers,
                                 data=json.dumps(query_payload),
                                 timeout=120, verify=False)
        response.raise_for_status()
        return response.json()
    except requests.exceptions.RequestException as e:
        if getattr(e, "response", None) is not None:
            print(f"  > 에러: Grafana API 호출 실패. 상태 코드: {e.response.status_code}, 응답: {e.response.text}")
        else:
            print(f"  > 에러: Grafana 데이터 API 호출 실패. {e}")
        return None



def get_all_placeholders(prs):
    placeholders = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    matches = re.findall(r'(\{\{.*?\}\})', p.text)
                    for match in matches:
                        placeholders.add(match)
    return list(placeholders)

def replace_text_in_presentation(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for p in shape.text_frame.paragraphs:
                if "{{" not in p.text: continue
                original_run = p.runs[0] if p.runs else None
                original_font = original_run.font if original_run else None
                temp_text = p.text
                text_changed = False
                for placeholder, value in replacements.items():
                    if placeholder in temp_text:
                        temp_text = temp_text.replace(placeholder, str(value))
                        text_changed = True
                if text_changed:
                    p.text = temp_text
                    if original_font:
                        for run in p.runs:
                            font = run.font
                            font.name = original_font.name
                            if original_font.size: font.size = original_font.size
                            if original_font.bold is not None: font.bold = original_font.bold
                            if original_font.italic is not None: font.italic = original_font.italic
                            try:
                                if original_font.color.rgb: font.color.rgb = RGBColor.from_string(str(original_font.color.rgb))
                            except: pass

# --- 메인 실행 로직 ---
if __name__ == "__main__":
    import urllib3
    from pptx.dml.color import RGBColor
    import time
    import gc
    urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)
    
    target_customer_path = None
    if len(sys.argv) > 1:
        target_customer_path = sys.argv[1]
        print(f"'{target_customer_path}'에 대한 개별 처리를 시작합니다.")
        all_files = find_templates_for_customer(BASE_TEMPLATE_DIR, target_customer_path)
    else:
        print("--- 전체 숫자/통계 삽입을 시작합니다 ---")
        all_files = find_all_templates(BASE_TEMPLATE_DIR)

    if not all_files:
        print("처리할 파일이 없습니다.")
        exit()
        
    date_info = calculate_previous_month_dates()
    start_ts = date_info['start_ts']
    end_ts = date_info['end_ts']
    
    dashboard_defs_cache = {}

    for i, file_path in enumerate(all_files):
        print("-" * 40)
        print(f"작업 {i+1}/{len(all_files)}: '{os.path.basename(file_path)}'")
        
        final_replacements = date_info["placeholders"].copy()
        relative_path = os.path.relpath(os.path.dirname(file_path), BASE_TEMPLATE_DIR)
        customer_name = relative_path.split(os.sep)[0]
        
        dashboard_def = dashboard_defs_cache.get(customer_name)
        if not dashboard_def:
            dashboard_uid = DASHBOARD_MAP.get(customer_name)
            if not dashboard_uid:
                print(f"  > 경고: DASHBOARD_MAP에 '{customer_name}' 정보 없음. 통계 삽입 건너뜁니다.")
            else:
                print(f"  > '{customer_name}' 대시보드(UID: {dashboard_uid}) 로딩 중...")
                dashboard_def = get_dashboard_definition(dashboard_uid)
                dashboard_defs_cache[customer_name] = dashboard_def
        
        try:
            prs = Presentation(file_path)
            
            grafana_failures = []
            if dashboard_def:
                all_panels_flat_list = dashboard_def.get('all_panels', [])
                all_ph = get_all_placeholders(prs)
                
                panel_data_cache = {}
                for ph in all_ph:
                    if ph in final_replacements: continue
                    try:
                        inner_text = ph.replace("{{", "").replace("}}", "")
                        parts = inner_text.rsplit('_', 1)
                        if len(parts) != 2: continue

                        panel_title_slug, query_letter = parts
                        
                        cache_key = (panel_title_slug, query_letter)
                        if cache_key not in panel_data_cache:
                            panel_title = panel_title_slug.replace("-", " ")
                            print(f"  > Grafana 조회: '{panel_title}' - '{query_letter}'쿼리")
                            panel = find_panel_by_title(all_panels_flat_list, panel_title)
                            if panel:
                                panel_data_cache[cache_key] = get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts)
                                time.sleep(1)

                            else:
                                print(f"  > 경고: '{panel_title}' 패널 없음.")
                                panel_data_cache[cache_key] = None
                        
                        stats_data = panel_data_cache.get(cache_key)
                        if not stats_data:
                            final_replacements[ph] = "N/A"
                            grafana_failures.append(ph)
                            continue

                        # 숫자 컬럼 자동 탐색
                        metric_values = None
                        frames = stats_data.get('results', {}).get(query_letter, {}).get('frames', [])
                        if frames:
                            fields = frames[0].get('schema', {}).get('fields', [])
                            for idx, field in enumerate(fields):
                                if field.get('type') == 'number':
                                    metric_values = frames[0].get('data', {}).get('values', [])[idx]
                                    break
                        
                        if metric_values:
                            valid_numbers = [v for v in metric_values if v is not None]
                            if valid_numbers:
                                max_val = max(valid_numbers)
                                mean_val = sum(valid_numbers) / len(valid_numbers)
                                final_replacements[ph] = SENTENCE_TEMPLATE.format(max=f"{max_val:.1f}", mean=f"{mean_val:.1f}")
                            else:
                                final_replacements[ph] = "N/A"
                                grafana_failures.append(ph)
                        else:
                            final_replacements[ph] = "N/A"
                            grafana_failures.append(ph)
                            
                    except (KeyError, IndexError, TypeError, ValueError):
                        final_replacements[ph] = "N/A"
                        grafana_failures.append(ph)

            replace_text_in_presentation(prs, final_replacements)
            
            if grafana_failures:
                print("\n  ⚠️  Grafana 데이터 조회에 실패한 플레이с홀더 목록:")
                for name in sorted(grafana_failures):
                    print(f"    - {name}")

            output_subdir = os.path.join(OUTPUT_DIR, relative_path)
            if not os.path.exists(output_subdir):
                os.makedirs(output_subdir)

            final_output_path = os.path.join(output_subdir, os.path.basename(file_path))
            prs.save(final_output_path)
            print(f"\n  > 최종 보고서 저장 완료: {final_output_path}")

            # 메모리 관리
            del prs
            if 'all_ph' in locals(): del all_ph
            if 'panel_data_cache' in locals(): del panel_data_cache
            gc.collect()

        except Exception as e:
            print(f"  > 에러 발생: {e}")
            continue

    print("\n--- 모든 숫자/통계 삽입 완료 ---")
