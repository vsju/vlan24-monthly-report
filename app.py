from flask import Flask, request, jsonify, send_file, send_from_directory, Response
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
import sys
import traceback
import zipfile
import io
import re
import shutil
import base64
import copy
import time
import hashlib
import json
import jwt
import logging
import urllib.parse
import requests
from datetime import datetime, timedelta

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import config
import image_processor
import stats_processor
import template_generator
import image_renderer

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)
CORS(app)

def validate_customer_name(name):
    """고객사 이름 유효성 검사 - Path traversal 방지"""
    if not name:
        return False, "이름이 비어있습니다."
    
    if '..' in name or '/' in name or '\\' in name:
        return False, "이름에 경로 구분자를 포함할 수 없습니다."
    
    if name.startswith('.'):
        return False, "이름이 .으로 시작할 수 없습니다."
    
    if any(c in name for c in ['<', '>', ':', '"', '|', '?', '*', '\x00']):
        return False, "이름에 특수문자를 포함할 수 없습니다."
    
    if name in ['template', 'completed_with_images', 'completed_final']:
        return False, "예약된 이름은 사용할 수 없습니다."
    
    return True, None

def validate_path_in_base(path, base_dir):
    """경로가 base_dir 내부인지 확인 - Path traversal 및 symlink 우회 방지"""
    real_path = os.path.realpath(path)
    real_base = os.path.realpath(base_dir)
    return real_path.startswith(real_base + os.sep) or real_path == real_base

def get_template_base_dir(template_type):
    """템플릿 타입에 따른 기본 디렉토리 반환"""
    if template_type == 'source':
        return config.BASE_TEMPLATE_DIR
    elif template_type == 'with_images':
        return config.OUTPUT_DIR_WITH_IMAGES
    elif template_type == 'final':
        return config.OUTPUT_DIR
    else:
        return config.BASE_TEMPLATE_DIR

def set_cell_text_preserve_format(cell, new_text):
    """셀 텍스트를 변경하면서 기존 서식(폰트, 크기, 색상, 정렬, 하이퍼링크 등) 모두 보존
    
    첫 번째 run의 text만 변경하고, 나머지 run과 paragraph는 그대로 유지
    이렇게 하면 멀티 run/paragraph 구조의 복잡한 서식도 보존됨
    """
    text_frame = cell.text_frame
    new_text = str(new_text)
    
    if len(text_frame.paragraphs) == 0:
        cell.text = new_text
        return
    
    first_para = text_frame.paragraphs[0]
    
    if len(first_para.runs) == 0:
        first_para.text = new_text
        return
    
    first_run = first_para.runs[0]
    first_run.text = new_text
    
    for run in list(first_para.runs)[1:]:
        run.text = ""

os.makedirs(config.BASE_TEMPLATE_DIR, exist_ok=True)
os.makedirs(config.BASE_IMAGE_DIR, exist_ok=True)
os.makedirs(config.OUTPUT_DIR_WITH_IMAGES, exist_ok=True)
os.makedirs(config.OUTPUT_DIR, exist_ok=True)

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({
        "status": "healthy",
        "grafana_url": config.GRAFANA_URL,
        "grafana_configured": bool(config.API_KEY),
        "directories": {
            "templates": os.path.exists(config.BASE_TEMPLATE_DIR),
            "images": os.path.exists(config.BASE_IMAGE_DIR),
            "output_images": os.path.exists(config.OUTPUT_DIR_WITH_IMAGES),
            "output_final": os.path.exists(config.OUTPUT_DIR)
        }
    })

@app.route('/api/config', methods=['GET'])
def get_config():
    return jsonify({
        "grafana_url": config.GRAFANA_URL,
        "dashboard_map": config.load_dashboard_map(),
        "directories": {
            "templates": config.BASE_TEMPLATE_DIR,
            "images": config.BASE_IMAGE_DIR,
            "output_with_images": config.OUTPUT_DIR_WITH_IMAGES,
            "output_final": config.OUTPUT_DIR
        }
    })

@app.route('/api/customers', methods=['GET'])
def list_customers():
    customers = []
    dashboard_map = config.load_dashboard_map()
    
    if os.path.exists(config.BASE_IMAGE_DIR):
        for item in os.listdir(config.BASE_IMAGE_DIR):
            item_path = os.path.join(config.BASE_IMAGE_DIR, item)
            if os.path.isdir(item_path) and item not in ['template', 'completed_with_images', 'completed_final']:
                customer_data = dashboard_map.get(item, {})
                if isinstance(customer_data, str):
                    customer_data = {"dashboard_uid": customer_data, "display_name": "", "contact_name": "", "contact_phone": "", "contact_email": ""}
                customers.append({
                    "name": item,
                    "display_name": customer_data.get('display_name', '') or '',
                    "dashboard_uid": customer_data.get('dashboard_uid', '') or '',
                    "contact_name": customer_data.get('contact_name', '') or '',
                    "contact_phone": customer_data.get('contact_phone', '') or '',
                    "contact_email": customer_data.get('contact_email', '') or '',
                    "has_images": len([f for f in os.listdir(item_path) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]) > 0
                })
    
    for customer_name, customer_data in dashboard_map.items():
        if not any(c["name"] == customer_name for c in customers):
            if isinstance(customer_data, str):
                customer_data = {"dashboard_uid": customer_data, "display_name": "", "contact_name": "", "contact_phone": "", "contact_email": ""}
            customers.append({
                "name": customer_name,
                "display_name": customer_data.get('display_name', '') or '',
                "dashboard_uid": customer_data.get('dashboard_uid', '') or '',
                "contact_name": customer_data.get('contact_name', '') or '',
                "contact_phone": customer_data.get('contact_phone', '') or '',
                "contact_email": customer_data.get('contact_email', '') or '',
                "has_images": False
            })
    
    return jsonify({"customers": sorted(customers, key=lambda x: x["name"])})

@app.route('/api/customers', methods=['POST'])
def create_customer():
    data = request.get_json() or {}
    customer_name = data.get('name', '').strip()
    display_name = data.get('display_name', '').strip()
    dashboard_uid = data.get('dashboard_uid', '').strip()
    contact_name = data.get('contact_name', '').strip()
    contact_phone = data.get('contact_phone', '').strip()
    contact_email = data.get('contact_email', '').strip()
    
    if not customer_name:
        return jsonify({"success": False, "error": "고객사 이름(폴더명)이 필요합니다."}), 400
    
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    customer_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
    template_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer_name)
    
    if not validate_path_in_base(customer_dir, config.BASE_IMAGE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    if not validate_path_in_base(template_dir, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    os.makedirs(customer_dir, exist_ok=True)
    os.makedirs(template_dir, exist_ok=True)
    
    config.set_customer_data(
        customer_name,
        dashboard_uid=dashboard_uid if dashboard_uid else None,
        display_name=display_name if display_name else None,
        contact_name=contact_name if contact_name else None,
        contact_phone=contact_phone if contact_phone else None,
        contact_email=contact_email if contact_email else None
    )
    
    return jsonify({
        "success": True,
        "message": f"고객사 '{customer_name}' 생성 완료",
        "customer": {
            "name": customer_name,
            "display_name": display_name,
            "dashboard_uid": dashboard_uid,
            "contact_name": contact_name,
            "contact_phone": contact_phone,
            "contact_email": contact_email,
            "directories": {
                "images": customer_dir,
                "templates": template_dir
            }
        }
    })

@app.route('/api/customers/<customer_name>', methods=['PUT'])
def update_customer(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    data = request.get_json() or {}
    dashboard_uid = data.get('dashboard_uid')
    display_name = data.get('display_name')
    contact_name = data.get('contact_name')
    contact_phone = data.get('contact_phone')
    contact_email = data.get('contact_email')
    
    updated_fields = []
    
    if dashboard_uid is not None:
        dashboard_uid = dashboard_uid.strip()
        if dashboard_uid:
            config.set_dashboard_uid(customer_name, dashboard_uid)
            updated_fields.append('dashboard_uid')
        else:
            config.delete_dashboard_mapping(customer_name)
    
    if display_name is not None:
        display_name = display_name.strip()
        config.set_display_name(customer_name, display_name)
        updated_fields.append('display_name')
    
    if contact_name is not None:
        config.set_customer_data(customer_name, contact_name=contact_name.strip())
        updated_fields.append('contact_name')
    
    if contact_phone is not None:
        config.set_customer_data(customer_name, contact_phone=contact_phone.strip())
        updated_fields.append('contact_phone')
    
    if contact_email is not None:
        config.set_customer_data(customer_name, contact_email=contact_email.strip())
        updated_fields.append('contact_email')
    
    return jsonify({
        "success": True,
        "message": f"고객사 '{customer_name}' 업데이트 완료",
        "updated_fields": updated_fields
    })

@app.route('/api/customers/<customer_name>', methods=['DELETE'])
def delete_customer(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    data = request.get_json() or {}
    delete_files = data.get('delete_files', False)
    
    config.delete_dashboard_mapping(customer_name)
    config.delete_customer_metadata(customer_name)
    
    if delete_files:
        customer_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
        template_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer_name)
        
        if not validate_path_in_base(customer_dir, config.BASE_IMAGE_DIR):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        if not validate_path_in_base(template_dir, config.BASE_TEMPLATE_DIR):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        
        deleted_dirs = []
        if os.path.exists(customer_dir):
            shutil.rmtree(customer_dir)
            deleted_dirs.append(customer_dir)
        if os.path.exists(template_dir):
            shutil.rmtree(template_dir)
            deleted_dirs.append(template_dir)
        
        return jsonify({
            "success": True,
            "message": f"고객사 '{customer_name}' 및 관련 파일 삭제 완료",
            "deleted_directories": deleted_dirs
        })
    else:
        return jsonify({
            "success": True,
            "message": f"고객사 '{customer_name}' 대시보드 매핑 삭제 완료 (파일은 유지)"
        })

@app.route('/api/customers/<customer_name>/subdirs', methods=['GET'])
def get_customer_subdirs(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    customer_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
    
    if not validate_path_in_base(customer_dir, config.BASE_IMAGE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    subdirs = []
    if os.path.exists(customer_dir) and os.path.isdir(customer_dir):
        for item in os.listdir(customer_dir):
            item_path = os.path.join(customer_dir, item)
            if os.path.isdir(item_path) and not item.startswith('.'):
                image_count = len([f for f in os.listdir(item_path) 
                                 if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))])
                subdirs.append({
                    "name": item,
                    "image_count": image_count
                })
    
    subdirs.sort(key=lambda x: x['name'])
    
    return jsonify({
        "success": True,
        "customer": customer_name,
        "subdirs": subdirs
    })

@app.route('/api/customers/<customer_name>/images', methods=['GET'])
def get_customer_images(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    subdir = request.args.get('subdir', '')
    
    if subdir:
        is_valid_subdir, subdir_error = validate_customer_name(subdir)
        if not is_valid_subdir:
            return jsonify({"success": False, "error": f"하위 디렉토리: {subdir_error}"}), 400
        target_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name, subdir)
    else:
        target_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
    
    if not validate_path_in_base(target_dir, config.BASE_IMAGE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    images = []
    if os.path.exists(target_dir) and os.path.isdir(target_dir):
        for item in os.listdir(target_dir):
            if item.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                item_path = os.path.join(target_dir, item)
                if os.path.isfile(item_path):
                    images.append({
                        "name": item,
                        "size": os.path.getsize(item_path)
                    })
    
    images.sort(key=lambda x: x['name'])
    
    return jsonify({
        "success": True,
        "customer": customer_name,
        "subdir": subdir,
        "images": images
    })

@app.route('/api/customers/<customer_name>/analyze-vms', methods=['GET'])
def analyze_customer_vms(customer_name):
    """고객사 이미지 폴더 분석: VM 목록 및 리소스 수 반환 (name 또는 display_name으로 검색)"""
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    customer_dir = os.path.join(config.BASE_IMAGE_DIR, customer_name)
    
    if not os.path.exists(customer_dir) or not os.path.isdir(customer_dir):
        found_name = config.find_customer_by_name_or_display(customer_name)
        if found_name:
            customer_dir = os.path.join(config.BASE_IMAGE_DIR, found_name)
            customer_name = found_name
    
    if not validate_path_in_base(customer_dir, config.BASE_IMAGE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(customer_dir) or not os.path.isdir(customer_dir):
        return jsonify({"success": False, "error": "고객사 폴더가 존재하지 않습니다."}), 404
    
    from template_generator import parse_vm_directory
    
    vms = []
    for item in sorted(os.listdir(customer_dir)):
        item_path = os.path.join(customer_dir, item)
        if os.path.isdir(item_path) and not item.startswith('.'):
            vm_name, vm_ip = parse_vm_directory(item)
            
            images = [f for f in os.listdir(item_path) 
                     if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]
            
            resources = []
            for img in sorted(images):
                img_name = os.path.splitext(img)[0]
                if '_:' in img_name:
                    parts = img_name.split('_:')
                    if len(parts) >= 2:
                        resource_part = parts[1]
                        if '_' in resource_part:
                            resource_name = resource_part.rsplit('_', 1)[0]
                        else:
                            resource_name = resource_part
                        resources.append({
                            "name": resource_name,
                            "filename": img_name,
                            "query": "A" if resource_name.upper() == "CPU" else "C"
                        })
            
            vms.append({
                "dir_name": item,
                "vm_name": vm_name,
                "ip": vm_ip,
                "image_count": len(images),
                "resources": resources,
                "pages_needed": (len(resources) + 1) // 2
            })
    
    return jsonify({
        "success": True,
        "customer": customer_name,
        "vms": vms,
        "total_vms": len(vms)
    })

@app.route('/api/dashboard-mapping', methods=['GET'])
def get_dashboard_mapping():
    return jsonify({
        "success": True,
        "mapping": config.load_dashboard_map()
    })

@app.route('/api/dashboard-mapping', methods=['PUT'])
def update_dashboard_mapping():
    data = request.get_json() or {}
    mapping = data.get('mapping', {})
    
    if not isinstance(mapping, dict):
        return jsonify({"success": False, "error": "mapping은 객체여야 합니다."}), 400
    
    config.save_dashboard_map(mapping)
    
    return jsonify({
        "success": True,
        "message": "대시보드 매핑 업데이트 완료",
        "mapping": mapping
    })

@app.route('/api/dashboard-mapping/<customer_name>', methods=['PUT'])
def update_single_dashboard_mapping(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    data = request.get_json() or {}
    dashboard_uid = data.get('dashboard_uid', '').strip()
    
    if dashboard_uid:
        config.set_dashboard_uid(customer_name, dashboard_uid)
        return jsonify({
            "success": True,
            "message": f"'{customer_name}' 대시보드 UID 설정: {dashboard_uid}"
        })
    else:
        config.delete_dashboard_mapping(customer_name)
        return jsonify({
            "success": True,
            "message": f"'{customer_name}' 대시보드 매핑 삭제"
        })

@app.route('/api/templates', methods=['GET'])
def list_templates():
    customer = request.args.get('customer', None)
    template_type = request.args.get('type', 'source')
    
    if template_type == 'source':
        base_dir = config.BASE_TEMPLATE_DIR
    elif template_type == 'with_images':
        base_dir = config.OUTPUT_DIR_WITH_IMAGES
    elif template_type == 'final':
        base_dir = config.OUTPUT_DIR
    else:
        base_dir = config.BASE_TEMPLATE_DIR
    
    templates = []
    
    if not os.path.exists(base_dir):
        return jsonify({
            "success": True,
            "templates": [],
            "count": 0,
            "base_dir": base_dir
        })
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
        
        search_dir = os.path.join(base_dir, customer)
        if not validate_path_in_base(search_dir, base_dir):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        
        if os.path.exists(search_dir):
            for root, dirs, files in os.walk(search_dir):
                for f in files:
                    if f.endswith('.pptx') and not f.startswith('~$'):
                        filepath = os.path.join(root, f)
                        rel_path = os.path.relpath(filepath, base_dir)
                        templates.append({
                            "filename": f,
                            "customer": customer,
                            "path": rel_path,
                            "size": os.path.getsize(filepath),
                            "modified": os.path.getmtime(filepath),
                            "type": template_type
                        })
    else:
        for root, dirs, files in os.walk(base_dir):
            for f in files:
                if f.endswith('.pptx') and not f.startswith('~$'):
                    filepath = os.path.join(root, f)
                    rel_path = os.path.relpath(filepath, base_dir)
                    customer_name = rel_path.split(os.sep)[0] if os.sep in rel_path else ""
                    templates.append({
                        "filename": f,
                        "customer": customer_name,
                        "path": rel_path,
                        "size": os.path.getsize(filepath),
                        "modified": os.path.getmtime(filepath),
                        "type": template_type
                    })
    
    return jsonify({
        "success": True,
        "templates": templates,
        "count": len(templates),
        "type": template_type
    })

@app.route('/api/templates/all', methods=['GET'])
def list_all_templates():
    """모든 유형의 템플릿을 한번에 반환"""
    def get_templates_from_dir(base_dir, template_type):
        templates = []
        if not os.path.exists(base_dir):
            return templates
        for root, dirs, files in os.walk(base_dir):
            for f in files:
                if f.endswith('.pptx') and not f.startswith('~$'):
                    filepath = os.path.join(root, f)
                    rel_path = os.path.relpath(filepath, base_dir)
                    customer_name = rel_path.split(os.sep)[0] if os.sep in rel_path else ""
                    templates.append({
                        "filename": f,
                        "customer": customer_name,
                        "path": rel_path,
                        "size": os.path.getsize(filepath),
                        "modified": os.path.getmtime(filepath),
                        "type": template_type
                    })
        return templates
    
    source_templates = get_templates_from_dir(config.BASE_TEMPLATE_DIR, 'source')
    with_images_templates = get_templates_from_dir(config.OUTPUT_DIR_WITH_IMAGES, 'with_images')
    final_templates = get_templates_from_dir(config.OUTPUT_DIR, 'final')
    
    return jsonify({
        "success": True,
        "templates": source_templates,
        "with_images": with_images_templates,
        "final": final_templates,
        "counts": {
            "source": len(source_templates),
            "with_images": len(with_images_templates),
            "final": len(final_templates)
        }
    })

@app.route('/api/templates/copy', methods=['POST'])
def copy_template():
    data = request.get_json() or {}
    source_path = data.get('source_path', '').strip()
    target_customer = data.get('target_customer', '').strip()
    new_filename = data.get('new_filename', '').strip()
    source_type = data.get('source_type', 'source')
    
    if not source_path:
        return jsonify({"success": False, "error": "원본 템플릿 경로가 필요합니다."}), 400
    if not target_customer:
        return jsonify({"success": False, "error": "대상 고객사가 필요합니다."}), 400
    
    is_valid, error_msg = validate_customer_name(target_customer)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    if source_type == 'source':
        source_base = config.BASE_TEMPLATE_DIR
    elif source_type == 'with_images':
        source_base = config.OUTPUT_DIR_WITH_IMAGES
    elif source_type == 'final':
        source_base = config.OUTPUT_DIR
    else:
        source_base = config.BASE_TEMPLATE_DIR
    
    source_full_path = os.path.join(source_base, source_path)
    
    if not validate_path_in_base(source_full_path, source_base):
        return jsonify({"success": False, "error": "잘못된 원본 경로입니다."}), 400
    
    if not os.path.exists(source_full_path):
        return jsonify({"success": False, "error": "원본 템플릿을 찾을 수 없습니다."}), 404
    
    target_dir = os.path.join(config.BASE_TEMPLATE_DIR, target_customer)
    if not validate_path_in_base(target_dir, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 대상 경로입니다."}), 400
    
    if not os.path.exists(target_dir):
        os.makedirs(target_dir, exist_ok=True)
    
    if not new_filename:
        new_filename = os.path.basename(source_path)
    
    if not new_filename.endswith('.pptx'):
        new_filename += '.pptx'
    
    target_full_path = os.path.join(target_dir, new_filename)
    
    if os.path.exists(target_full_path):
        return jsonify({"success": False, "error": f"대상 파일이 이미 존재합니다: {new_filename}"}), 400
    
    shutil.copy2(source_full_path, target_full_path)
    
    return jsonify({
        "success": True,
        "message": f"템플릿 복사 완료: {target_customer}/{new_filename}",
        "target_path": f"{target_customer}/{new_filename}"
    })

@app.route('/api/templates/<path:template_path>', methods=['DELETE'])
def delete_template(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    os.remove(full_path)
    
    return jsonify({
        "success": True,
        "message": f"템플릿 '{template_path}' 삭제 완료"
    })

@app.route('/api/templates/<path:template_path>/info', methods=['GET'])
def get_template_info(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        prs = Presentation(full_path)
        
        slides_info = []
        for idx, slide in enumerate(prs.slides):
            shapes_info = []
            placeholders = []
            
            for shape in slide.shapes:
                shape_info = {
                    "id": shape.shape_id,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                
                if shape.has_text_frame:
                    text = ""
                    for p in shape.text_frame.paragraphs:
                        text += p.text + "\n"
                    shape_info["text"] = text.strip()
                    
                    matches = re.findall(r'(\{\{.*?\}\})', text)
                    placeholders.extend(matches)
                
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_text = ""
                    table_cells = []
                    for row_idx, row in enumerate(shape.table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = ""
                            if cell.text_frame:
                                for p in cell.text_frame.paragraphs:
                                    cell_text += p.text + "\n"
                            cell_text = cell_text.strip()
                            if cell_text:
                                table_text += cell_text + "\n"
                                table_cells.append({
                                    "row": row_idx,
                                    "col": col_idx,
                                    "text": cell_text
                                })
                    shape_info["text"] = table_text.strip()
                    shape_info["table_cells"] = table_cells
                    shape_info["table_rows"] = len(shape.table.rows)
                    shape_info["table_cols"] = len(shape.table.columns)
                    
                    matches = re.findall(r'(\{\{.*?\}\})', table_text)
                    placeholders.extend(matches)
                
                shapes_info.append(shape_info)
            
            slides_info.append({
                "index": idx,
                "shapes": shapes_info,
                "shape_count": len(shapes_info),
                "placeholders": list(set(placeholders))
            })
        
        return jsonify({
            "success": True,
            "template": {
                "path": template_path,
                "slide_count": len(prs.slides),
                "slides": slides_info
            }
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/slides/<int:slide_index>/preview', methods=['GET'])
def get_slide_preview(template_path, slide_index):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        from pptx.util import Emu
        import subprocess
        import tempfile
        
        temp_dir = tempfile.mkdtemp()
        temp_pdf = os.path.join(temp_dir, "output.pdf")
        temp_png = os.path.join(temp_dir, f"output-{slide_index + 1}.png")
        
        try:
            result = subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', temp_dir, full_path
            ], capture_output=True, timeout=60)
            
            pdf_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(full_path))[0] + ".pdf")
            
            if os.path.exists(pdf_file):
                result = subprocess.run([
                    'pdftoppm', '-png', '-f', str(slide_index + 1), '-l', str(slide_index + 1),
                    pdf_file, os.path.join(temp_dir, "output")
                ], capture_output=True, timeout=30)
                
                png_files = [f for f in os.listdir(temp_dir) if f.startswith("output") and f.endswith(".png")]
                if png_files:
                    with open(os.path.join(temp_dir, png_files[0]), 'rb') as f:
                        image_data = base64.b64encode(f.read()).decode('utf-8')
                    
                    shutil.rmtree(temp_dir)
                    return jsonify({
                        "success": True,
                        "image": f"data:image/png;base64,{image_data}"
                    })
        except Exception as e:
            pass
        finally:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
        
        prs = Presentation(full_path)
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        shapes_preview = []
        
        for shape in slide.shapes:
            shape_data = {
                "name": shape.name,
                "type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            if shape.has_text_frame:
                text = ""
                for p in shape.text_frame.paragraphs:
                    text += p.text + "\n"
                shape_data["text"] = text.strip()
            shapes_preview.append(shape_data)
        
        return jsonify({
            "success": True,
            "preview_type": "shapes",
            "shapes": shapes_preview,
            "message": "이미지 미리보기가 불가능하여 Shape 정보를 반환합니다."
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/previews', methods=['GET'])
def get_all_slide_previews(template_path):
    """전체 슬라이드 미리보기를 일괄 생성하여 반환"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        import subprocess
        import tempfile
        
        prs = Presentation(full_path)
        slide_count = len(prs.slides)
        previews = []
        
        temp_dir = tempfile.mkdtemp()
        pdf_generated = False
        
        try:
            result = subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', temp_dir, full_path
            ], capture_output=True, timeout=120)
            
            pdf_file = os.path.join(temp_dir, os.path.splitext(os.path.basename(full_path))[0] + ".pdf")
            
            if os.path.exists(pdf_file):
                result = subprocess.run([
                    'pdftoppm', '-png', pdf_file, os.path.join(temp_dir, "slide")
                ], capture_output=True, timeout=60)
                pdf_generated = True
                
                for i in range(slide_count):
                    png_pattern = f"slide-{i+1}.png"
                    png_files = [f for f in os.listdir(temp_dir) if f == png_pattern or f == f"slide-{str(i+1).zfill(2)}.png" or f == f"slide-{str(i+1).zfill(3)}.png"]
                    
                    if png_files:
                        with open(os.path.join(temp_dir, png_files[0]), 'rb') as f:
                            image_data = base64.b64encode(f.read()).decode('utf-8')
                        previews.append({
                            "index": i,
                            "type": "image",
                            "image": f"data:image/png;base64,{image_data}"
                        })
                    else:
                        previews.append({
                            "index": i,
                            "type": "text",
                            "shapes": get_slide_shapes_info(prs.slides[i])
                        })
        except Exception as e:
            pdf_generated = False
        finally:
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
        
        if not pdf_generated or len(previews) == 0:
            previews = []
            for i, slide in enumerate(prs.slides):
                previews.append({
                    "index": i,
                    "type": "text",
                    "shapes": get_slide_shapes_info(slide)
                })
        
        return jsonify({
            "success": True,
            "slide_count": slide_count,
            "previews": previews
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

def get_slide_shapes_info(slide):
    """슬라이드의 Shape 정보를 추출"""
    shapes_info = []
    for shape in slide.shapes:
        shape_data = {
            "name": shape.name,
            "type": str(shape.shape_type).replace("MSO_SHAPE_TYPE.", ""),
            "id": shape.shape_id
        }
        if shape.has_text_frame:
            text = ""
            for p in shape.text_frame.paragraphs:
                text += p.text + "\n"
            shape_data["text"] = text.strip()
        shapes_info.append(shape_data)
    return shapes_info

@app.route('/api/templates/<path:template_path>/shapes/<int:slide_index>/<int:shape_id>', methods=['PUT'])
def update_shape(template_path, slide_index, shape_id):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    new_name = data.get('name')
    new_text = data.get('text')
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        target_shape = None
        
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                target_shape = shape
                break
        
        if not target_shape:
            return jsonify({"success": False, "error": "Shape를 찾을 수 없습니다."}), 404
        
        changes = []
        
        if new_name is not None:
            old_name = target_shape.name
            target_shape.name = new_name
            changes.append(f"이름: '{old_name}' -> '{new_name}'")
        
        if new_text is not None:
            if target_shape.has_text_frame:
                old_text = ""
                for p in target_shape.text_frame.paragraphs:
                    old_text += p.text
                
                if target_shape.text_frame.paragraphs:
                    target_shape.text_frame.paragraphs[0].runs[0].text = new_text if target_shape.text_frame.paragraphs[0].runs else new_text
                changes.append(f"텍스트 변경")
            elif target_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                row_idx = data.get('row')
                col_idx = data.get('col')
                if row_idx is not None and col_idx is not None:
                    table = target_shape.table
                    if row_idx < len(table.rows) and col_idx < len(table.columns):
                        cell = table.cell(row_idx, col_idx)
                        tf = cell.text_frame
                        if tf and tf.paragraphs:
                            for para in tf.paragraphs:
                                for run in para.runs:
                                    run.text = ""
                            if tf.paragraphs[0].runs:
                                tf.paragraphs[0].runs[0].text = new_text
                            else:
                                tf.paragraphs[0].text = new_text
                        changes.append(f"테이블 셀 ({row_idx}, {col_idx}) 텍스트 변경")
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": "Shape 업데이트 완료",
            "changes": changes
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/placeholders', methods=['GET'])
def get_placeholders(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        prs = Presentation(full_path)
        
        all_placeholders = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        matches = re.findall(r'(\{\{.*?\}\})', p.text)
                        for match in matches:
                            all_placeholders.append({
                                "placeholder": match,
                                "slide_index": slide_idx,
                                "shape_id": shape.shape_id,
                                "shape_name": shape.name
                            })
        
        unique_placeholders = list(set([p["placeholder"] for p in all_placeholders]))
        
        return jsonify({
            "success": True,
            "placeholders": all_placeholders,
            "unique_placeholders": unique_placeholders,
            "count": len(unique_placeholders)
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/placeholders', methods=['POST'])
def add_placeholder(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    slide_index = data.get('slide_index', 0)
    shape_id = data.get('shape_id')
    placeholder = data.get('placeholder', '')
    
    if not placeholder:
        return jsonify({"success": False, "error": "placeholder가 필요합니다."}), 400
    
    if not placeholder.startswith('{{') or not placeholder.endswith('}}'):
        placeholder = '{{' + placeholder + '}}'
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        
        if shape_id:
            for shape in slide.shapes:
                if shape.shape_id == shape_id and shape.has_text_frame:
                    if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                        shape.text_frame.paragraphs[0].runs[0].text += " " + placeholder
                    else:
                        shape.text_frame.text += " " + placeholder
                    break
        else:
            from pptx.util import Inches
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(0.5))
            textbox.text_frame.text = placeholder
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"플레이스홀더 '{placeholder}' 추가 완료"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/placeholders', methods=['DELETE'])
def delete_placeholder(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    placeholder = data.get('placeholder', '')
    
    if not placeholder:
        return jsonify({"success": False, "error": "placeholder가 필요합니다."}), 400
    
    try:
        prs = Presentation(full_path)
        removed_count = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, '')
                                removed_count += 1
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"플레이스홀더 '{placeholder}' 삭제 완료",
            "removed_count": removed_count
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/slides/<int:slide_index>', methods=['DELETE'])
def delete_slide(template_path, slide_index):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        if len(prs.slides) <= 1:
            return jsonify({"success": False, "error": "마지막 슬라이드는 삭제할 수 없습니다."}), 400
        
        slide_id = prs.slides._sldIdLst[slide_index].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[slide_index]
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"슬라이드 {slide_index + 1} 삭제 완료",
            "remaining_slides": len(prs.slides)
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/slides/reorder', methods=['POST'])
def reorder_slides(template_path):
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    old_index = data.get('old_index')
    new_index = data.get('new_index')
    
    if old_index is None or new_index is None:
        return jsonify({"success": False, "error": "old_index와 new_index가 필요합니다."}), 400
    
    try:
        prs = Presentation(full_path)
        
        if old_index >= len(prs.slides) or new_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide_elem = prs.slides._sldIdLst[old_index]
        prs.slides._sldIdLst.remove(slide_elem)
        prs.slides._sldIdLst.insert(new_index, slide_elem)
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"슬라이드 순서 변경 완료: {old_index + 1} -> {new_index + 1}"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/slides/duplicate', methods=['POST'])
def duplicate_slide(template_path):
    """슬라이드 복제 API - 순번 자동 증가 + VM 정보 교체 + VM 목록 표에 행 추가"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    slide_index = data.get('slide_index')
    vm_name = data.get('vm_name', '')
    vm_ip = data.get('vm_ip', '')
    vm_os = data.get('vm_os', '')
    
    if slide_index is None:
        return jsonify({"success": False, "error": "slide_index가 필요합니다."}), 400
    
    if not vm_name:
        return jsonify({"success": False, "error": "vm_name이 필요합니다."}), 400
    
    try:
        from copy import deepcopy
        from lxml import etree
        
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        source_slide = prs.slides[slide_index]
        
        title_pattern = re.compile(r'^(\d+)\.(\d+)\s+(.+?)\s*\((.+?)\)$')
        current_major = None
        current_minor = 0
        title_shape = None
        
        for shape in source_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                match = title_pattern.match(text)
                if match:
                    current_major = int(match.group(1))
                    current_minor = int(match.group(2))
                    title_shape = shape
                    break
        
        if current_major is None:
            for shape in source_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    simple_pattern = re.compile(r'^(\d+)\.(\d+)\s+')
                    match = simple_pattern.match(text)
                    if match:
                        current_major = int(match.group(1))
                        current_minor = int(match.group(2))
                        title_shape = shape
                        break
        
        if current_major is None:
            current_major = 3
            current_minor = 0
        
        max_minor = current_minor
        seq_pattern = re.compile(rf'^{current_major}\.(\d+)')
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    match = seq_pattern.match(text)
                    if match:
                        minor = int(match.group(1))
                        if minor > max_minor:
                            max_minor = minor
        
        new_minor = max_minor + 1
        new_sequence = f"{current_major}.{new_minor}"
        
        slide_layout = source_slide.slide_layout
        new_slide = prs.slides.add_slide(slide_layout)
        
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        
        for shape in source_slide.shapes:
            el = shape._element
            new_el = deepcopy(el)
            new_slide.shapes._spTree.append(new_el)
        
        new_title = f"{new_sequence} {vm_name}"
        if vm_ip:
            new_title += f" ({vm_ip})"
        
        for shape in new_slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if title_pattern.match(text) or (current_major and seq_pattern.match(text)):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if title_pattern.match(run.text.strip()) or seq_pattern.match(run.text.strip()):
                                run.text = new_title
                                break
                        else:
                            continue
                        break
                    else:
                        if shape.text_frame.paragraphs:
                            shape.text_frame.paragraphs[0].runs[0].text = new_title
                    break
        
        new_slide_elem = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.remove(new_slide_elem)
        prs.slides._sldIdLst.insert(slide_index + 1, new_slide_elem)
        
        vm_table_updated = False
        table_pattern = re.compile(r'VM\s*명|서버\s*명', re.IGNORECASE)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    if table.rows and len(table.rows) > 0:
                        first_row_text = ' '.join([cell.text for cell in table.rows[0].cells])
                        if table_pattern.search(first_row_text):
                            new_row_idx = len(table.rows)
                            table.add_row()
                            new_row = table.rows[new_row_idx]
                            
                            col_count = len(table.columns)
                            if col_count >= 1:
                                new_row.cells[0].text = str(new_row_idx)
                            if col_count >= 2:
                                new_row.cells[1].text = vm_name
                            if col_count >= 3:
                                new_row.cells[2].text = vm_ip
                            if col_count >= 4:
                                new_row.cells[3].text = vm_os
                            
                            vm_table_updated = True
                            break
            if vm_table_updated:
                break
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"슬라이드 복제 완료: {new_sequence} {vm_name}",
            "new_sequence": new_sequence,
            "new_slide_index": slide_index + 1,
            "vm_table_updated": vm_table_updated
        })
        
    except Exception as e:
        import traceback
        return jsonify({"success": False, "error": str(e), "traceback": traceback.format_exc()}), 500

@app.route('/api/templates/<path:template_path>/placeholders/update', methods=['PUT'])
def update_placeholder(template_path):
    """플레이스홀더 수정 API - 기존 플레이스홀더를 새 텍스트로 교체"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    old_placeholder = data.get('old_placeholder', '')
    new_placeholder = data.get('new_placeholder', '')
    
    if not old_placeholder:
        return jsonify({"success": False, "error": "old_placeholder가 필요합니다."}), 400
    
    try:
        prs = Presentation(full_path)
        replaced_count = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_placeholder in run.text:
                                run.text = run.text.replace(old_placeholder, new_placeholder)
                                replaced_count += 1
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"플레이스홀더 수정 완료: '{old_placeholder}' -> '{new_placeholder}'",
            "replaced_count": replaced_count
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/tables', methods=['GET'])
def get_tables(template_path):
    """템플릿 내 표 목록 조회"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        prs = Presentation(full_path)
        tables = []
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    rows_data = []
                    headers = []
                    
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        if row_idx == 0:
                            headers = row_data
                        rows_data.append(row_data)
                    
                    tables.append({
                        "slide_index": slide_idx,
                        "shape_id": shape.shape_id,
                        "shape_name": shape.name,
                        "row_count": len(table.rows),
                        "col_count": len(table.columns),
                        "headers": headers,
                        "rows": rows_data
                    })
        
        return jsonify({
            "success": True,
            "tables": tables,
            "count": len(tables)
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/tables/<int:slide_index>/<int:shape_id>/rows', methods=['POST'])
def add_table_row(template_path, slide_index, shape_id):
    """표에 행 추가"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    row_data = data.get('row_data', [])
    insert_at = data.get('insert_at', None)
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        target_table = None
        target_shape = None
        
        for shape in slide.shapes:
            if shape.has_table and shape.shape_id == shape_id:
                target_table = shape.table
                target_shape = shape
                break
        
        if not target_table or not target_shape:
            return jsonify({"success": False, "error": "표를 찾을 수 없습니다."}), 404
        
        tbl = target_shape._element.graphic.graphicData.tbl
        tr_elements = tbl.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tr')
        
        if len(tr_elements) < 1:
            return jsonify({"success": False, "error": "복제할 행이 없습니다."}), 400
        
        if insert_at is not None and len(tr_elements) > 1:
            clamped_idx = max(1, min(insert_at, len(tr_elements) - 1))
            template_tr = tr_elements[clamped_idx]
        else:
            template_tr = tr_elements[-1]
        new_tr = copy.deepcopy(template_tr)
        
        tc_elements = new_tr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tc')
        for col_idx, tc in enumerate(tc_elements):
            t_elements = tc.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}t')
            for t_idx, t in enumerate(t_elements):
                if t_idx == 0:
                    if col_idx < len(row_data):
                        t.text = str(row_data[col_idx])
                    else:
                        t.text = ""
                else:
                    t.text = ""
        
        if insert_at is not None and 0 < insert_at < len(tr_elements):
            reference_tr = tr_elements[insert_at]
            reference_tr.addprevious(new_tr)
            new_row_idx = insert_at
        else:
            tbl.append(new_tr)
            new_row_idx = len(tr_elements)
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"행 추가 완료 (행 {new_row_idx + 1})",
            "row_index": new_row_idx
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/tables/<int:slide_index>/<int:shape_id>', methods=['PUT'])
def update_table_all_rows(template_path, slide_index, shape_id):
    """표 전체 행 수정"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    all_rows = data.get('rows', [])
    
    if not all_rows:
        return jsonify({"success": False, "error": "저장할 데이터가 없습니다."}), 400
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        target_table = None
        
        for shape in slide.shapes:
            if shape.has_table and shape.shape_id == shape_id:
                target_table = shape.table
                break
        
        if not target_table:
            return jsonify({"success": False, "error": "표를 찾을 수 없습니다."}), 404
        
        updated_count = 0
        skip_header = data.get('skip_header', True)
        
        for row_index, row_data in enumerate(all_rows):
            if skip_header and row_index == 0:
                continue
            if row_index >= len(target_table.rows):
                break
            row = target_table.rows[row_index]
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(target_table.columns):
                    set_cell_text_preserve_format(row.cells[col_idx], cell_text)
            updated_count += 1
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"전체 저장 완료 ({updated_count}개 행)"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/tables/<int:slide_index>/<int:shape_id>/rows/<int:row_index>', methods=['PUT'])
def update_table_row(template_path, slide_index, shape_id, row_index):
    """표 행 수정"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    data = request.get_json() or {}
    row_data = data.get('row_data', [])
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        target_table = None
        
        for shape in slide.shapes:
            if shape.has_table and shape.shape_id == shape_id:
                target_table = shape.table
                break
        
        if not target_table:
            return jsonify({"success": False, "error": "표를 찾을 수 없습니다."}), 404
        
        if row_index >= len(target_table.rows):
            return jsonify({"success": False, "error": "행 인덱스가 범위를 벗어났습니다."}), 400
        
        row = target_table.rows[row_index]
        for col_idx, cell_text in enumerate(row_data):
            if col_idx < len(target_table.columns):
                set_cell_text_preserve_format(row.cells[col_idx], cell_text)
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"행 수정 완료 (행 {row_index + 1})"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/templates/<path:template_path>/tables/<int:slide_index>/<int:shape_id>/rows/<int:row_index>', methods=['DELETE'])
def delete_table_row(template_path, slide_index, shape_id, row_index):
    """표 행 삭제"""
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "템플릿을 찾을 수 없습니다."}), 404
    
    try:
        prs = Presentation(full_path)
        
        if slide_index >= len(prs.slides):
            return jsonify({"success": False, "error": "슬라이드 인덱스가 범위를 벗어났습니다."}), 400
        
        slide = prs.slides[slide_index]
        target_shape = None
        
        for shape in slide.shapes:
            if shape.has_table and shape.shape_id == shape_id:
                target_shape = shape
                break
        
        if not target_shape:
            return jsonify({"success": False, "error": "표를 찾을 수 없습니다."}), 404
        
        table = target_shape.table
        
        if row_index >= len(table.rows):
            return jsonify({"success": False, "error": "행 인덱스가 범위를 벗어났습니다."}), 400
        
        if row_index == 0:
            return jsonify({"success": False, "error": "헤더 행은 삭제할 수 없습니다."}), 400
        
        if len(table.rows) <= 2:
            return jsonify({"success": False, "error": "최소 2개의 행이 필요합니다."}), 400
        
        tbl = target_shape._element.graphic.graphicData.tbl
        tr_elements = tbl.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tr')
        if row_index < len(tr_elements):
            tr_elements[row_index].getparent().remove(tr_elements[row_index])
        
        prs.save(full_path)
        
        return jsonify({
            "success": True,
            "message": f"행 삭제 완료 (행 {row_index + 1})"
        })
        
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/api/process/images', methods=['POST'])
def process_images():
    data = request.get_json() or {}
    customer_name = data.get('customer_name', None)
    
    if customer_name:
        is_valid, error_msg = validate_customer_name(customer_name)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
    
    results = image_processor.process_images(customer_name)
    
    return jsonify(results), 200 if results["success"] else 400

@app.route('/api/process/statistics', methods=['POST'])
def process_statistics():
    data = request.get_json() or {}
    customer_name = data.get('customer_name', None)
    
    if customer_name:
        is_valid, error_msg = validate_customer_name(customer_name)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
    
    results = stats_processor.process_statistics(customer_name)
    
    return jsonify(results), 200 if results["success"] else 400

@app.route('/api/templates/generate', methods=['POST'])
def generate_template():
    """마스터 템플릿 기반 고객사별 템플릿 자동 생성"""
    data = request.get_json() or {}
    customer_name = data.get('customer_name')
    master_template = data.get('master_template')
    
    if not customer_name:
        return jsonify({"success": False, "error": "customer_name이 필요합니다."}), 400
    
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    if not master_template:
        master_files = [f for f in os.listdir(config.BASE_TEMPLATE_DIR) 
                       if f.endswith('.pptx') and '{{NAME}}' in f]
        if master_files:
            master_template = os.path.join(config.BASE_TEMPLATE_DIR, master_files[0])
        else:
            all_pptx = [f for f in os.listdir(config.BASE_TEMPLATE_DIR) if f.endswith('.pptx')]
            if all_pptx:
                master_template = os.path.join(config.BASE_TEMPLATE_DIR, all_pptx[0])
            else:
                return jsonify({"success": False, "error": "마스터 템플릿을 찾을 수 없습니다."}), 400
    else:
        if not os.path.isabs(master_template):
            master_template = os.path.join(config.BASE_TEMPLATE_DIR, master_template)
    
    if not validate_path_in_base(master_template, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 템플릿 경로입니다."}), 400
    
    if not os.path.exists(master_template):
        return jsonify({"success": False, "error": f"마스터 템플릿이 존재하지 않습니다: {master_template}"}), 404
    
    display_name = config.get_display_name(customer_name)
    
    results = template_generator.generate_customer_template(
        master_template, 
        customer_name,
        display_name=display_name
    )
    
    return jsonify(results), 200 if results["success"] else 400


@app.route('/api/templates/<customer_name>/add-vm', methods=['POST'])
def add_vm_to_template_api(customer_name):
    """기존 템플릿에 새 VM 슬라이드 추가"""
    data = request.get_json() or {}
    vm_dir_name = data.get('vm_dir_name')
    template_filename = data.get('template_filename')
    master_template = data.get('master_template')
    seq_number = data.get('seq_number')
    
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    if not vm_dir_name:
        return jsonify({"success": False, "error": "vm_dir_name이 필요합니다."}), 400
    
    if '..' in vm_dir_name or vm_dir_name.startswith('/'):
        return jsonify({"success": False, "error": "잘못된 VM 폴더명입니다."}), 400
    
    if template_filename:
        template_path = os.path.join(config.BASE_TEMPLATE_DIR, customer_name, template_filename)
    else:
        customer_template_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer_name)
        if os.path.exists(customer_template_dir):
            pptx_files = [f for f in os.listdir(customer_template_dir) if f.endswith('.pptx')]
            if pptx_files:
                template_path = os.path.join(customer_template_dir, pptx_files[0])
            else:
                return jsonify({"success": False, "error": "고객사 템플릿이 없습니다."}), 404
        else:
            return jsonify({"success": False, "error": "고객사 템플릿 폴더가 없습니다."}), 404
    
    if not validate_path_in_base(template_path, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 템플릿 경로입니다."}), 400
    
    if not os.path.exists(template_path):
        return jsonify({"success": False, "error": f"템플릿이 존재하지 않습니다: {template_path}"}), 404
    
    if not master_template:
        master_files = [f for f in os.listdir(config.BASE_TEMPLATE_DIR) 
                       if f.endswith('.pptx') and os.path.isfile(os.path.join(config.BASE_TEMPLATE_DIR, f))]
        if master_files:
            master_template = os.path.join(config.BASE_TEMPLATE_DIR, master_files[0])
        else:
            return jsonify({"success": False, "error": "마스터 템플릿을 찾을 수 없습니다."}), 400
    else:
        if not os.path.isabs(master_template):
            master_template = os.path.join(config.BASE_TEMPLATE_DIR, master_template)
    
    if not validate_path_in_base(master_template, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 마스터 템플릿 경로입니다."}), 400
    
    if not os.path.exists(master_template):
        return jsonify({"success": False, "error": f"마스터 템플릿이 존재하지 않습니다: {master_template}"}), 404
    
    results = template_generator.add_vm_to_template(
        template_path=template_path,
        master_template_path=master_template,
        vm_dir_name=vm_dir_name,
        customer_name=customer_name,
        seq_number=seq_number
    )
    
    return jsonify(results), 200 if results["success"] else 400


@app.route('/api/process/all', methods=['POST'])
def process_all():
    results = {
        "success": True,
        "images": {},
        "statistics": {},
        "errors": []
    }
    
    try:
        results["images"] = image_processor.process_images()
        
        if not results["images"]["success"]:
            results["success"] = False
            results["errors"].append("이미지 삽입 실패")
        else:
            results["statistics"] = stats_processor.process_statistics()
            
            if not results["statistics"]["success"]:
                results["success"] = False
                results["errors"].append("통계 삽입 실패")
    
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"전체 프로세스 실패: {str(e)}")
    
    return jsonify(results), 200 if results["success"] else 400

@app.route('/api/upload/templates', methods=['POST'])
def upload_templates():
    if 'files' not in request.files:
        return jsonify({"success": False, "error": "파일이 없습니다."}), 400
    
    files = request.files.getlist('files')
    customer = request.form.get('customer', '')
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
    
    uploaded_files = []
    errors = []
    
    for file in files:
        if file.filename == '':
            continue
        
        if file and file.filename.endswith('.pptx'):
            filename = secure_filename(file.filename)
            
            if customer:
                target_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer)
                if not validate_path_in_base(target_dir, config.BASE_TEMPLATE_DIR):
                    return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
            else:
                target_dir = config.BASE_TEMPLATE_DIR
            
            os.makedirs(target_dir, exist_ok=True)
            filepath = os.path.join(target_dir, filename)
            file.save(filepath)
            uploaded_files.append(filepath)
        else:
            errors.append(f"'{file.filename}'은(는) .pptx 파일이 아닙니다.")
    
    return jsonify({
        "success": len(uploaded_files) > 0,
        "uploaded_files": uploaded_files,
        "errors": errors
    })

@app.route('/api/report-root/subdirs', methods=['GET'])
def list_report_root_subdirs():
    subdirs = []
    
    if os.path.exists(config.BASE_IMAGE_DIR):
        for item in os.listdir(config.BASE_IMAGE_DIR):
            item_path = os.path.join(config.BASE_IMAGE_DIR, item)
            if os.path.isdir(item_path) and item not in ['template', 'completed_with_images', 'completed_final']:
                image_count = len([f for f in os.listdir(item_path) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))])
                subdirs.append({
                    "name": item,
                    "image_count": image_count
                })
    
    return jsonify({"subdirs": sorted(subdirs, key=lambda x: x['name'])})


@app.route('/api/upload/images', methods=['POST'])
def upload_images():
    if 'files' not in request.files:
        return jsonify({"success": False, "error": "파일이 없습니다."}), 400
    
    files = request.files.getlist('files')
    customer = request.form.get('customer', '')
    subdir = request.form.get('subdir', '')
    root_report = request.form.get('root_report', 'false') == 'true'
    
    if not customer and not root_report:
        return jsonify({"success": False, "error": "customer가 필요합니다."}), 400
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
    
    if subdir:
        is_valid_subdir, subdir_error = validate_customer_name(subdir)
        if not is_valid_subdir:
            return jsonify({"success": False, "error": f"하위 디렉토리: {subdir_error}"}), 400
    
    uploaded_files = []
    errors = []
    
    if root_report:
        if subdir:
            target_dir = os.path.join(config.BASE_IMAGE_DIR, subdir)
        else:
            target_dir = config.BASE_IMAGE_DIR
    elif subdir:
        target_dir = os.path.join(config.BASE_IMAGE_DIR, customer, subdir)
    else:
        target_dir = os.path.join(config.BASE_IMAGE_DIR, customer)
    
    if not validate_path_in_base(target_dir, config.BASE_IMAGE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    os.makedirs(target_dir, exist_ok=True)
    
    for file in files:
        if file.filename == '':
            continue
        
        if file and file.filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
            filename = secure_filename(file.filename)
            filepath = os.path.join(target_dir, filename)
            file.save(filepath)
            uploaded_files.append(filepath)
        else:
            errors.append(f"'{file.filename}'은(는) 지원되지 않는 이미지 형식입니다.")
    
    return jsonify({
        "success": len(uploaded_files) > 0,
        "uploaded_files": uploaded_files,
        "errors": errors
    })

@app.route('/api/download/results', methods=['GET'])
def download_results():
    customer = request.args.get('customer', None)
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
        
        target_dir = os.path.join(config.OUTPUT_DIR, customer)
        if not validate_path_in_base(target_dir, config.OUTPUT_DIR):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        
        if not os.path.exists(target_dir):
            return jsonify({"success": False, "error": f"고객사 '{customer}' 결과가 없습니다."}), 404
        
        memory_file = io.BytesIO()
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(target_dir):
                for file in files:
                    if file.endswith('.pptx'):
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, config.OUTPUT_DIR)
                        zipf.write(file_path, arcname)
        
        memory_file.seek(0)
        return send_file(
            memory_file,
            mimetype='application/zip',
            as_attachment=True,
            download_name=f'{customer}_results.zip'
        )
    else:
        memory_file = io.BytesIO()
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(config.OUTPUT_DIR):
                for file in files:
                    if file.endswith('.pptx'):
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, config.OUTPUT_DIR)
                        zipf.write(file_path, arcname)
        
        memory_file.seek(0)
        return send_file(
            memory_file,
            mimetype='application/zip',
            as_attachment=True,
            download_name='all_results.zip'
        )

@app.route('/api/download/templates', methods=['GET'])
def download_templates():
    template_type = request.args.get('type', 'final')
    customer = request.args.get('customer', None)
    
    type_to_dir = {
        'source': config.BASE_TEMPLATE_DIR,
        'with_images': config.OUTPUT_DIR_WITH_IMAGES,
        'final': config.OUTPUT_DIR
    }
    
    if template_type not in type_to_dir:
        return jsonify({"success": False, "error": f"잘못된 유형입니다. (source, with_images, final 중 선택)"}), 400
    
    base_dir = type_to_dir[template_type]
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
        
        target_dir = os.path.join(base_dir, customer)
        if not validate_path_in_base(target_dir, base_dir):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        
        if not os.path.exists(target_dir):
            return jsonify({"success": False, "error": f"'{customer}' 고객사의 {template_type} 템플릿이 없습니다."}), 404
        
        memory_file = io.BytesIO()
        file_count = 0
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(target_dir):
                for file in files:
                    if file.endswith('.pptx'):
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, base_dir)
                        zipf.write(file_path, arcname)
                        file_count += 1
        
        if file_count == 0:
            return jsonify({"success": False, "error": "다운로드할 템플릿이 없습니다."}), 404
        
        memory_file.seek(0)
        return send_file(
            memory_file,
            mimetype='application/zip',
            as_attachment=True,
            download_name=f'{customer}_{template_type}_templates.zip'
        )
    else:
        memory_file = io.BytesIO()
        file_count = 0
        with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(base_dir):
                for file in files:
                    if file.endswith('.pptx'):
                        file_path = os.path.join(root, file)
                        arcname = os.path.relpath(file_path, base_dir)
                        zipf.write(file_path, arcname)
                        file_count += 1
        
        if file_count == 0:
            return jsonify({"success": False, "error": "다운로드할 템플릿이 없습니다."}), 404
        
        memory_file.seek(0)
        return send_file(
            memory_file,
            mimetype='application/zip',
            as_attachment=True,
            download_name=f'all_{template_type}_templates.zip'
        )

@app.route('/api/files/list', methods=['GET'])
def list_files():
    file_type = request.args.get('type', 'results')
    customer = request.args.get('customer', None)
    
    if file_type == 'results':
        base_dir = config.OUTPUT_DIR
    elif file_type == 'templates':
        base_dir = config.BASE_TEMPLATE_DIR
    else:
        return jsonify({"success": False, "error": "잘못된 타입입니다."}), 400
    
    files_list = []
    
    if customer:
        is_valid, error_msg = validate_customer_name(customer)
        if not is_valid:
            return jsonify({"success": False, "error": error_msg}), 400
        
        target_dir = os.path.join(base_dir, customer)
        if not validate_path_in_base(target_dir, base_dir):
            return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
        
        if os.path.exists(target_dir):
            for root, dirs, files in os.walk(target_dir):
                for file in files:
                    if file.endswith('.pptx'):
                        rel_path = os.path.relpath(os.path.join(root, file), base_dir)
                        files_list.append(rel_path)
    else:
        for root, dirs, files in os.walk(base_dir):
            for file in files:
                if file.endswith('.pptx'):
                    rel_path = os.path.relpath(os.path.join(root, file), base_dir)
                    files_list.append(rel_path)
    
    return jsonify({
        "success": True,
        "files": files_list,
        "count": len(files_list)
    })

WOPI_JWT_SECRET = os.getenv('WOPI_JWT_SECRET')
if not WOPI_JWT_SECRET:
    import warnings
    warnings.warn("WOPI_JWT_SECRET not set! Using insecure default. Set WOPI_JWT_SECRET environment variable.")
    WOPI_JWT_SECRET = 'XSPalFKvv/MyVxXZoeR99IKDg4O0Mj8j3DIf7Omwt0A='
ONLYOFFICE_URL = os.getenv('ONLYOFFICE_URL', 'http://192.168.10.77:8080')
ONLYOFFICE_EXTERNAL_URL = os.getenv('ONLYOFFICE_EXTERNAL_URL', 'http://121.78.82.22:8502')
WOPI_BASE_URL = os.getenv('WOPI_BASE_URL', 'http://192.168.10.30:5001')

def generate_wopi_token(file_id, user_id="admin", expires_hours=24):
    """WOPI 액세스 토큰 생성"""
    payload = {
        "file_id": file_id,
        "user_id": user_id,
        "exp": datetime.utcnow() + timedelta(hours=expires_hours),
        "iat": datetime.utcnow()
    }
    return jwt.encode(payload, WOPI_JWT_SECRET, algorithm="HS256")

def verify_wopi_token(token, file_id):
    """WOPI 액세스 토큰 검증"""
    try:
        payload = jwt.decode(token, WOPI_JWT_SECRET, algorithms=["HS256"])
        return payload.get("file_id") == file_id
    except jwt.ExpiredSignatureError:
        return False
    except jwt.InvalidTokenError:
        return False

def verify_onlyoffice_jwt(token, file_id):
    """OnlyOffice가 보내는 JWT 토큰 검증 (Authorization 헤더용)
    
    OnlyOffice config token의 document.url에 file_id가 포함되어 있는지 확인하여
    다른 파일에 대한 토큰 재사용 공격을 방지합니다.
    URL 인코딩된 file_id도 비교합니다 (한글 파일명 지원).
    """
    try:
        payload = jwt.decode(token, WOPI_JWT_SECRET, algorithms=["HS256"])
        
        document_url = urllib.parse.unquote(payload.get('document', {}).get('url', ''))
        decoded_file_id = urllib.parse.unquote(file_id)
        
        if decoded_file_id in document_url or file_id in document_url:
            return True
        
        wopi_src = urllib.parse.unquote(payload.get('wopi_src', ''))
        if decoded_file_id in wopi_src or file_id in wopi_src:
            return True
        
        app.logger.debug(f"JWT verification failed: file_id={file_id}, document_url={document_url}, wopi_src={wopi_src}")
        return False
    except jwt.ExpiredSignatureError:
        app.logger.warning(f"JWT token expired for file_id={file_id}")
        return False
    except jwt.InvalidTokenError as e:
        app.logger.warning(f"Invalid JWT token for file_id={file_id}: {e}")
        return False

def get_token_from_request(file_id):
    """요청에서 토큰 추출 및 검증 (query param 또는 Authorization 헤더)"""
    access_token = request.args.get('access_token', '')
    
    if access_token and verify_wopi_token(access_token, file_id):
        return True
    
    auth_header = request.headers.get('Authorization', '')
    if auth_header.startswith('Bearer '):
        bearer_token = auth_header[7:]
        if verify_onlyoffice_jwt(bearer_token, file_id):
            return True
    
    return False

def get_file_path_from_id(file_id):
    """파일 ID에서 실제 경로 추출"""
    decoded = base64.urlsafe_b64decode(file_id.encode()).decode()
    return decoded

def get_file_id_from_path(file_path):
    """파일 경로에서 ID 생성"""
    return base64.urlsafe_b64encode(file_path.encode()).decode()

wopi_locks = {}

def get_file_version(full_path):
    """파일 버전 해시 생성"""
    if os.path.exists(full_path):
        stat = os.stat(full_path)
        return hashlib.sha256(f"{stat.st_mtime}_{stat.st_size}".encode()).hexdigest()[:8]
    return "0"

@app.route('/wopi/files/<file_id>', methods=['GET'])
def wopi_check_file_info(file_id):
    """WOPI CheckFileInfo - 파일 메타데이터 반환"""
    app.logger.info(f"WOPI CheckFileInfo request: file_id={file_id}")
    
    if not get_token_from_request(file_id):
        app.logger.warning(f"WOPI CheckFileInfo: Invalid token for file_id={file_id}")
        return jsonify({"error": "Invalid token"}), 401
    
    try:
        file_path = get_file_path_from_id(file_id)
        app.logger.debug(f"WOPI CheckFileInfo: file_path={file_path}")
        full_path = os.path.join(config.BASE_TEMPLATE_DIR, file_path)
        
        if not validate_path_in_base(full_path, config.BASE_TEMPLATE_DIR):
            return jsonify({"error": "Invalid path"}), 400
        
        if not os.path.exists(full_path):
            return jsonify({"error": "File not found"}), 404
        
        stat = os.stat(full_path)
        file_hash = hashlib.sha256(open(full_path, 'rb').read()).hexdigest()
        
        version = get_file_version(full_path)
        
        return jsonify({
            "BaseFileName": os.path.basename(file_path),
            "Size": stat.st_size,
            "OwnerId": "admin",
            "UserId": "admin",
            "UserFriendlyName": "Admin User",
            "Version": version,
            "LastModifiedTime": datetime.fromtimestamp(stat.st_mtime).isoformat() + "Z",
            "UserCanWrite": True,
            "UserCanNotWriteRelative": True,
            "SupportsUpdate": True,
            "SupportsLocks": True,
            "SupportsCobalt": False,
            "SupportsGetLock": True,
            "SupportsFolders": False,
            "SupportsDeleteFile": False,
            "SupportsRename": False,
            "SupportsUserInfo": False
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/wopi/files/<file_id>', methods=['POST'])
def wopi_lock_operations(file_id):
    """WOPI Lock/Unlock/RefreshLock/GetLock 처리"""
    x_wopi_override = request.headers.get('X-WOPI-Override', '').upper()
    x_wopi_lock = request.headers.get('X-WOPI-Lock', '')
    x_wopi_old_lock = request.headers.get('X-WOPI-OldLock', '')
    
    app.logger.info(f"WOPI Lock operation: {x_wopi_override}, file_id={file_id}, lock={x_wopi_lock}")
    
    if not get_token_from_request(file_id):
        app.logger.warning(f"WOPI Lock: Invalid token for file_id={file_id}")
        return jsonify({"error": "Invalid token"}), 401
    
    try:
        file_path = get_file_path_from_id(file_id)
        full_path = os.path.join(config.BASE_TEMPLATE_DIR, file_path)
        
        if not validate_path_in_base(full_path, config.BASE_TEMPLATE_DIR):
            return jsonify({"error": "Invalid path"}), 400
        
        current_lock = wopi_locks.get(file_id, '')
        version = get_file_version(full_path)
        
        if x_wopi_override == 'GET_LOCK':
            response = Response(status=200)
            response.headers['X-WOPI-Lock'] = current_lock
            response.headers['X-WOPI-ItemVersion'] = version
            return response
        
        elif x_wopi_override == 'LOCK':
            if current_lock and current_lock != x_wopi_lock:
                response = Response(status=409)
                response.headers['X-WOPI-Lock'] = current_lock
                response.headers['X-WOPI-LockFailureReason'] = 'File is locked by another user'
                return response
            
            wopi_locks[file_id] = x_wopi_lock
            app.logger.info(f"WOPI Lock acquired: file_id={file_id}, lock={x_wopi_lock}")
            
            response = Response(status=200)
            response.headers['X-WOPI-ItemVersion'] = version
            return response
        
        elif x_wopi_override == 'UNLOCK':
            if current_lock and current_lock != x_wopi_lock:
                response = Response(status=409)
                response.headers['X-WOPI-Lock'] = current_lock
                return response
            
            wopi_locks.pop(file_id, None)
            app.logger.info(f"WOPI Lock released: file_id={file_id}")
            
            response = Response(status=200)
            response.headers['X-WOPI-ItemVersion'] = version
            return response
        
        elif x_wopi_override == 'REFRESH_LOCK':
            if current_lock and current_lock != x_wopi_lock:
                response = Response(status=409)
                response.headers['X-WOPI-Lock'] = current_lock
                return response
            
            wopi_locks[file_id] = x_wopi_lock
            app.logger.info(f"WOPI Lock refreshed: file_id={file_id}")
            
            response = Response(status=200)
            response.headers['X-WOPI-ItemVersion'] = version
            return response
        
        elif x_wopi_override == 'UNLOCK_AND_RELOCK':
            if current_lock and current_lock != x_wopi_old_lock:
                response = Response(status=409)
                response.headers['X-WOPI-Lock'] = current_lock
                return response
            
            wopi_locks[file_id] = x_wopi_lock
            app.logger.info(f"WOPI Lock reacquired: file_id={file_id}")
            
            response = Response(status=200)
            response.headers['X-WOPI-ItemVersion'] = version
            return response
        
        else:
            return jsonify({"error": f"Unknown operation: {x_wopi_override}"}), 400
    
    except Exception as e:
        app.logger.error(f"WOPI Lock error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/wopi/files/<file_id>/contents', methods=['GET'])
def wopi_get_file(file_id):
    """WOPI GetFile - 파일 다운로드"""
    app.logger.info(f"WOPI GetFile request: file_id={file_id}")
    
    if not get_token_from_request(file_id):
        app.logger.warning(f"WOPI GetFile: Invalid token for file_id={file_id}")
        return jsonify({"error": "Invalid token"}), 401
    
    try:
        file_path = get_file_path_from_id(file_id)
        app.logger.debug(f"WOPI GetFile: file_path={file_path}")
        full_path = os.path.join(config.BASE_TEMPLATE_DIR, file_path)
        
        if not validate_path_in_base(full_path, config.BASE_TEMPLATE_DIR):
            return jsonify({"error": "Invalid path"}), 400
        
        if not os.path.exists(full_path):
            return jsonify({"error": "File not found"}), 404
        
        version = get_file_version(full_path)
        
        response = send_file(
            full_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=False,
            download_name=os.path.basename(file_path)
        )
        response.headers['X-WOPI-ItemVersion'] = version
        return response
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/wopi/files/<file_id>/contents', methods=['POST'])
def wopi_put_file(file_id):
    """WOPI PutFile - 파일 저장"""
    x_wopi_lock = request.headers.get('X-WOPI-Lock', '')
    app.logger.info(f"WOPI PutFile request: file_id={file_id}, content_length={request.content_length}, lock={x_wopi_lock}")
    
    if not get_token_from_request(file_id):
        app.logger.warning(f"WOPI PutFile: Invalid token for file_id={file_id}")
        return jsonify({"error": "Invalid token"}), 401
    
    try:
        file_path = get_file_path_from_id(file_id)
        app.logger.debug(f"WOPI PutFile: file_path={file_path}")
        full_path = os.path.join(config.BASE_TEMPLATE_DIR, file_path)
        
        if not validate_path_in_base(full_path, config.BASE_TEMPLATE_DIR):
            return jsonify({"error": "Invalid path"}), 400
        
        current_lock = wopi_locks.get(file_id, '')
        if current_lock and current_lock != x_wopi_lock:
            app.logger.warning(f"WOPI PutFile: Lock mismatch for file_id={file_id}")
            response = Response(status=409)
            response.headers['X-WOPI-Lock'] = current_lock
            return response
        
        backup_path = full_path + ".bak"
        if os.path.exists(full_path):
            shutil.copy2(full_path, backup_path)
        
        with open(full_path, 'wb') as f:
            f.write(request.get_data())
        
        if os.path.exists(backup_path):
            os.remove(backup_path)
        
        version = get_file_version(full_path)
        app.logger.info(f"WOPI PutFile success: file_id={file_id}, new_version={version}")
        
        response = Response(status=200)
        response.headers['X-WOPI-ItemVersion'] = version
        return response
        
    except Exception as e:
        app.logger.error(f"WOPI PutFile error: {e}")
        if 'backup_path' in locals() and os.path.exists(backup_path):
            shutil.move(backup_path, full_path)
        return jsonify({"error": str(e)}), 500

def validate_onlyoffice_jwt(req):
    """OnlyOffice JWT 토큰 검증 (Authorization: Bearer 헤더)"""
    auth_header = req.headers.get('Authorization', '')
    if auth_header.startswith('Bearer '):
        token = auth_header[7:]
        try:
            jwt.decode(token, WOPI_JWT_SECRET, algorithms=["HS256"])
            return True, None
        except jwt.ExpiredSignatureError:
            return False, "Token expired"
        except jwt.InvalidTokenError as e:
            return False, f"Invalid token: {str(e)}"
    return True, None

@app.route('/api/onlyoffice/download/<path:template_path>', methods=['GET'])
def onlyoffice_download(template_path):
    """OnlyOffice에서 파일 다운로드용 엔드포인트"""
    is_valid, error = validate_onlyoffice_jwt(request)
    if not is_valid:
        app.logger.warning(f"OnlyOffice download JWT validation failed: {error}")
    
    template_type = request.args.get('type', 'source')
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return jsonify({"error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"error": "파일을 찾을 수 없습니다."}), 404
    
    app.logger.info(f"OnlyOffice download: {template_path}")
    
    return send_file(
        full_path,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        as_attachment=False,
        download_name=os.path.basename(template_path)
    )

DOCUMENT_KEY_MAPPING_FILE = os.path.join(os.path.dirname(__file__), 'document_key_mapping.json')
DOCUMENT_KEY_LOCK_FILE = os.path.join(os.path.dirname(__file__), 'document_key_mapping.lock')

import fcntl

def _ensure_lock_file_exists():
    """lock 파일이 존재하는지 확인하고 없으면 생성"""
    if not os.path.exists(DOCUMENT_KEY_LOCK_FILE):
        lock_dir = os.path.dirname(DOCUMENT_KEY_LOCK_FILE)
        if lock_dir and not os.path.exists(lock_dir):
            os.makedirs(lock_dir, exist_ok=True)
        with open(DOCUMENT_KEY_LOCK_FILE, 'w') as f:
            pass

def set_document_key(key, template_path):
    """document_key_mapping에 키 추가 및 저장 (전체 과정에 exclusive lock 적용)"""
    try:
        _ensure_lock_file_exists()
        
        with open(DOCUMENT_KEY_LOCK_FILE, 'r+') as lock_file:
            fcntl.flock(lock_file.fileno(), fcntl.LOCK_EX)
            try:
                mapping = {}
                if os.path.exists(DOCUMENT_KEY_MAPPING_FILE):
                    try:
                        with open(DOCUMENT_KEY_MAPPING_FILE, 'r') as f:
                            data = json.load(f)
                            mapping = data if isinstance(data, dict) else {}
                    except (json.JSONDecodeError, IOError):
                        mapping = {}
                
                mapping[key] = template_path
                
                if len(mapping) > 1000:
                    keys_to_remove = list(mapping.keys())[:len(mapping) - 500]
                    for k in keys_to_remove:
                        del mapping[k]
                
                with open(DOCUMENT_KEY_MAPPING_FILE, 'w') as f:
                    json.dump(mapping, f, ensure_ascii=False, indent=2)
            finally:
                fcntl.flock(lock_file.fileno(), fcntl.LOCK_UN)
    except Exception as e:
        app.logger.warning(f"Failed to set document_key: {e}")

def get_template_path_by_key(key):
    """document_key로 template_path 조회 (shared lock 적용)"""
    try:
        if not os.path.exists(DOCUMENT_KEY_MAPPING_FILE):
            return None
        
        _ensure_lock_file_exists()
        
        with open(DOCUMENT_KEY_LOCK_FILE, 'r') as lock_file:
            fcntl.flock(lock_file.fileno(), fcntl.LOCK_SH)
            try:
                with open(DOCUMENT_KEY_MAPPING_FILE, 'r') as f:
                    data = json.load(f)
                    mapping = data if isinstance(data, dict) else {}
                    return mapping.get(key)
            except (json.JSONDecodeError, IOError):
                return None
            finally:
                fcntl.flock(lock_file.fileno(), fcntl.LOCK_UN)
    except Exception as e:
        app.logger.warning(f"Failed to get template_path_by_key: {e}")
        return None

@app.route('/api/onlyoffice/callback', methods=['POST'])
def onlyoffice_callback():
    """OnlyOffice 콜백 - 문서 저장 처리"""
    try:
        is_valid, error = validate_onlyoffice_jwt(request)
        if not is_valid:
            app.logger.warning(f"OnlyOffice callback JWT validation failed: {error}")
        
        data = request.json
        app.logger.info(f"OnlyOffice callback: {data}")
        
        status = data.get('status')
        
        if status == 2 or status == 6:
            download_url = data.get('url')
            file_key = data.get('key', '')
            
            full_template_key = get_template_path_by_key(file_key)
            if not full_template_key:
                app.logger.error(f"OnlyOffice callback: Unknown key - {file_key}")
                return jsonify({"error": 1})
            
            if ':' in full_template_key:
                template_type, template_path = full_template_key.split(':', 1)
            else:
                template_type = 'source'
                template_path = full_template_key
            
            base_dir = get_template_base_dir(template_type)
            full_path = os.path.join(base_dir, template_path)
            
            if not validate_path_in_base(full_path, base_dir):
                return jsonify({"error": 1})
            
            if download_url:
                try:
                    # 외부 URL을 내부 URL로 변환 (Zabbix 서버는 외부 IP로 접근 불가)
                    internal_download_url = download_url.replace(
                        ONLYOFFICE_EXTERNAL_URL.rstrip('/'),
                        ONLYOFFICE_URL.rstrip('/')
                    )
                    app.logger.info(f"OnlyOffice callback: Downloading from {internal_download_url} (original: {download_url})")
                    response = requests.get(internal_download_url, timeout=30)
                    if response.ok:
                        backup_path = full_path + ".bak"
                        if os.path.exists(full_path):
                            shutil.copy2(full_path, backup_path)
                        
                        with open(full_path, 'wb') as f:
                            f.write(response.content)
                        
                        if os.path.exists(backup_path):
                            os.remove(backup_path)
                        
                        app.logger.info(f"OnlyOffice callback: File saved - {template_path}")
                    else:
                        app.logger.error(f"OnlyOffice callback: Download failed - status={response.status_code}, url={download_url}")
                        return jsonify({"error": 1})
                except requests.exceptions.RequestException as e:
                    app.logger.error(f"OnlyOffice callback: Download exception - {e}")
                    return jsonify({"error": 1})
            else:
                app.logger.error(f"OnlyOffice callback: No download URL provided")
                return jsonify({"error": 1})
        
        return jsonify({"error": 0})
        
    except Exception as e:
        app.logger.error(f"OnlyOffice callback error: {e}")
        return jsonify({"error": 1})

BACKEND_INTERNAL_URL = os.getenv('BACKEND_INTERNAL_URL', 'http://192.168.10.30:5001')

@app.route('/api/onlyoffice/editor-config', methods=['GET'])
def get_editor_config():
    """OnlyOffice 에디터 설정 반환 (DocsAPI.DocEditor 방식)"""
    template_path = request.args.get('template')
    
    if not template_path:
        return jsonify({"success": False, "error": "템플릿 경로가 필요합니다."}), 400
    
    full_path = os.path.join(config.BASE_TEMPLATE_DIR, template_path)
    
    if not validate_path_in_base(full_path, config.BASE_TEMPLATE_DIR):
        return jsonify({"success": False, "error": "잘못된 경로입니다."}), 400
    
    if not os.path.exists(full_path):
        return jsonify({"success": False, "error": "파일을 찾을 수 없습니다."}), 404
    
    file_id = get_file_id_from_path(template_path)
    
    document_key = hashlib.sha256(
        f"{file_id}_{os.path.getmtime(full_path)}".encode()
    ).hexdigest()[:20]
    
    set_document_key(document_key, template_path)
    
    backend_base = BACKEND_INTERNAL_URL.rstrip('/')
    encoded_path = urllib.parse.quote(template_path, safe='')
    
    document_url = f"{backend_base}/api/onlyoffice/download/{encoded_path}"
    callback_url = f"{backend_base}/api/onlyoffice/callback"
    
    editor_config = {
        "document": {
            "fileType": "pptx",
            "key": document_key,
            "title": os.path.basename(template_path),
            "url": document_url
        },
        "documentType": "slide",
        "editorConfig": {
            "callbackUrl": callback_url,
            "lang": "ko",
            "mode": "edit",
            "user": {
                "id": "admin",
                "name": "Admin"
            },
            "customization": {
                "autosave": True,
                "forcesave": True,
                "hideRightMenu": False,
                "hideRulers": False,
                "compactHeader": False
            }
        }
    }
    
    jwt_token = jwt.encode(editor_config, WOPI_JWT_SECRET, algorithm="HS256")
    editor_config["token"] = jwt_token
    
    return jsonify({
        "success": True,
        "config": editor_config,
        "onlyoffice_url": ONLYOFFICE_EXTERNAL_URL
    })

@app.route('/api/onlyoffice/editor-page', methods=['GET'])
def get_editor_page():
    """OnlyOffice 에디터 HTML 페이지 반환 (sandbox 제한 우회)"""
    template_path = request.args.get('template')
    template_type = request.args.get('type', 'source')
    initial_slide = request.args.get('slide', '0')
    refresh_ts = request.args.get('refresh', '')
    try:
        initial_slide = int(initial_slide)
    except ValueError:
        initial_slide = 0
    
    if not template_path:
        return "템플릿 경로가 필요합니다.", 400
    
    base_dir = get_template_base_dir(template_type)
    full_path = os.path.join(base_dir, template_path)
    
    if not validate_path_in_base(full_path, base_dir):
        return "잘못된 경로입니다.", 400
    
    if not os.path.exists(full_path):
        return "파일을 찾을 수 없습니다.", 404
    
    file_id = get_file_id_from_path(template_path)
    
    key_base = f"{template_type}_{file_id}_{os.path.getmtime(full_path)}"
    if refresh_ts:
        key_base += f"_{refresh_ts}"
    
    document_key = hashlib.sha256(key_base.encode()).hexdigest()[:20]
    
    full_template_key = f"{template_type}:{template_path}"
    set_document_key(document_key, full_template_key)
    
    backend_base = BACKEND_INTERNAL_URL.rstrip('/')
    encoded_path = urllib.parse.quote(template_path, safe='')
    
    document_url = f"{backend_base}/api/onlyoffice/download/{encoded_path}?type={template_type}"
    callback_url = f"{backend_base}/api/onlyoffice/callback?type={template_type}"
    
    editor_config = {
        "document": {
            "fileType": "pptx",
            "key": document_key,
            "title": os.path.basename(template_path),
            "url": document_url
        },
        "documentType": "slide",
        "editorConfig": {
            "callbackUrl": callback_url,
            "lang": "ko",
            "mode": "edit",
            "user": {
                "id": "admin",
                "name": "Admin"
            },
            "customization": {
                "autosave": True,
                "forcesave": True,
                "hideRightMenu": False,
                "hideRulers": False,
                "compactHeader": False
            }
        }
    }
    
    jwt_token = jwt.encode(editor_config, WOPI_JWT_SECRET, algorithm="HS256")
    editor_config["token"] = jwt_token
    
    config_json = json.dumps(editor_config)
    
    html_page = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>OnlyOffice Editor - {os.path.basename(template_path)}</title>
    <style>
        html, body {{
            margin: 0;
            padding: 0;
            height: 100%;
            overflow: hidden;
        }}
        #onlyoffice-editor {{
            width: 100%;
            height: 100%;
        }}
    </style>
</head>
<body>
    <div id="onlyoffice-editor"></div>
    <script type="text/javascript" src="{ONLYOFFICE_EXTERNAL_URL}/web-apps/apps/api/documents/api.js"></script>
    <script type="text/javascript">
        var config = {config_json};
        
        var initialSlide = {initial_slide};
        
        function initEditor() {{
            if (typeof DocsAPI !== 'undefined') {{
                window.docEditor = new DocsAPI.DocEditor("onlyoffice-editor", config);
                console.log("OnlyOffice Editor initialized successfully");
                
                if (initialSlide > 0) {{
                    setTimeout(function() {{
                        try {{
                            var pageNum = initialSlide + 1;
                            window.docEditor.setCurrentPage(pageNum);
                            console.log("Navigated to slide (1-based):", pageNum);
                        }} catch (err) {{
                            console.error("Error navigating to initial slide:", err);
                        }}
                    }}, 1500);
                }}
            }} else {{
                console.error("DocsAPI not loaded, retrying...");
                setTimeout(initEditor, 500);
            }}
        }}
        
        window.addEventListener('message', function(event) {{
            if (event.data && event.data.action === 'goToSlide') {{
                var slideIndex = event.data.slideIndex || 0;
                var pageNum = slideIndex + 1;
                console.log("Received goToSlide message for slide (0-based):", slideIndex, "-> page (1-based):", pageNum);
                if (window.docEditor) {{
                    try {{
                        window.docEditor.setCurrentPage(pageNum);
                        console.log("Navigated to slide:", pageNum);
                    }} catch (err) {{
                        console.error("Error navigating to slide:", err);
                    }}
                }}
            }}
        }});
        
        if (document.readyState === 'complete') {{
            initEditor();
        }} else {{
            window.onload = initEditor;
        }}
    </script>
</body>
</html>"""
    
    return html_page, 200, {'Content-Type': 'text/html; charset=utf-8'}

@app.route('/api/render/panels/<customer_name>', methods=['GET'])
def get_render_panels(customer_name):
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    result = image_renderer.get_panels_for_customer(customer_name)
    return jsonify(result)


@app.route('/api/render/all', methods=['POST'])
def render_all_images():
    data = request.get_json() or {}
    customer_name = data.get('customer_name', '').strip()
    
    if not customer_name:
        return jsonify({"success": False, "error": "고객사 이름이 필요합니다."}), 400
    
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    logs = []
    def log_func(msg):
        logs.append(msg)
    
    result = image_renderer.render_all_panels(customer_name, log_func=log_func)
    result["logs"] = logs
    return jsonify(result)


@app.route('/api/render/selected', methods=['POST'])
def render_selected_images():
    data = request.get_json() or {}
    customer_name = data.get('customer_name', '').strip()
    panel_ids = data.get('panel_ids', [])
    
    if not customer_name:
        return jsonify({"success": False, "error": "고객사 이름이 필요합니다."}), 400
    
    if not panel_ids:
        return jsonify({"success": False, "error": "렌더링할 패널을 선택해주세요."}), 400
    
    is_valid, error_msg = validate_customer_name(customer_name)
    if not is_valid:
        return jsonify({"success": False, "error": error_msg}), 400
    
    logs = []
    def log_func(msg):
        logs.append(msg)
    
    result = image_renderer.render_selected_panels(customer_name, panel_ids, log_func=log_func)
    result["logs"] = logs
    return jsonify(result)


@app.route('/api/render/all-customers', methods=['POST'])
def render_all_customers():
    logs = []
    def log_func(msg):
        logs.append(msg)
        app.logger.info(msg)
    
    result = image_renderer.render_all_customers(log_func=log_func)
    result["logs"] = logs
    return jsonify(result)


@app.errorhandler(Exception)
def handle_error(error):
    app.logger.error(f"Unhandled exception: {str(error)}\n{traceback.format_exc()}")
    return jsonify({
        "error": str(error),
        "traceback": traceback.format_exc()
    }), 500

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5001))
    app.run(host='0.0.0.0', port=port, debug=True)
