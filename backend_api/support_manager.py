import os
import datetime
from pptx import Presentation
from pptx.util import Pt
import config

def add_support_history(file_path, support_entries):
    results = {
        "success": True,
        "added_entries": 0,
        "errors": []
    }
    
    try:
        if not os.path.exists(file_path):
            results["success"] = False
            results["errors"].append(f"파일을 찾을 수 없음: {file_path}")
            return results
        
        prs = Presentation(file_path)
        
        if len(prs.slides) < 2:
            results["success"] = False
            results["errors"].append("슬라이드가 충분하지 않습니다 (최소 2개 필요).")
            return results
        
        target_slide = prs.slides[-2]
        
        table_shape = None
        for shape in target_slide.shapes:
            if shape.has_table:
                table_shape = shape
                break
        
        if not table_shape:
            results["success"] = False
            results["errors"].append("지원 내역 표를 찾을 수 없습니다.")
            return results
        
        table = table_shape.table
        
        for entry in support_entries:
            date = entry.get('date', datetime.date.today().strftime("%Y-%m-%d"))
            content = entry.get('content', '')
            
            if not content:
                continue
            
            row = table.add_row()
            row.cells[0].text = date
            row.cells[1].text = content
            
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
            
            results["added_entries"] += 1
        
        prs.save(file_path)
        
        if results["added_entries"] == 0:
            results["success"] = False
            results["errors"].append("추가된 항목이 없습니다.")
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"지원 내역 추가 실패: {str(e)}")
    
    return results
