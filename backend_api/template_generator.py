import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import config

def get_dashboard_panels(dashboard_uid):
    url = f"{config.GRAFANA_URL.rstrip('/')}/api/dashboards/uid/{dashboard_uid}"
    headers = {"Authorization": f"Bearer {config.API_KEY}"}
    
    try:
        response = requests.get(url, headers=headers, timeout=20, verify=config.VERIFY_SSL)
        response.raise_for_status()
        dashboard = response.json().get('dashboard', {})
        panels = []
        
        def extract_panels(panel_list):
            for panel in panel_list:
                if panel.get('type') != 'row':
                    panels.append(panel)
                if panel.get('type') == 'row' and 'panels' in panel:
                    extract_panels(panel['panels'])
        
        extract_panels(dashboard.get('panels', []))
        return panels
    except Exception as e:
        raise Exception(f"Grafana 대시보드 조회 실패: {str(e)}")

def generate_template(customer_name, selected_panels=None):
    results = {
        "success": True,
        "generated_file": None,
        "placeholders": [],
        "errors": []
    }
    
    try:
        if customer_name not in config.DASHBOARD_MAP:
            results["success"] = False
            results["errors"].append(f"'{customer_name}' 대시보드 매핑이 없습니다.")
            return results
        
        dashboard_uid = config.DASHBOARD_MAP[customer_name]
        panels = get_dashboard_panels(dashboard_uid)
        
        if not panels:
            results["success"] = False
            results["errors"].append("패널을 찾을 수 없습니다.")
            return results
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{customer_name} 운영 보고서"
        subtitle.text = "{{DATE_RANGE}}"
        
        for idx, panel in enumerate(panels, 1):
            panel_title = panel.get('title', f'패널{idx}')
            targets = panel.get('targets', [])
            
            if selected_panels and panel_title not in selected_panels:
                continue
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = panel_title
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            chart_placeholder = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5),
                Inches(9), Inches(4)
            )
            chart_frame = chart_placeholder.text_frame
            chart_frame.text = f"[{panel_title} 차트 위치]"
            chart_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            for query_idx, target in enumerate(targets):
                ref_id = target.get('refId', chr(65 + query_idx))
                placeholder = f"{{{{{panel_title}_{ref_id}}}}}"
                results["placeholders"].append({
                    "panel": panel_title,
                    "query": ref_id,
                    "placeholder": placeholder
                })
                
                stats_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(5.8 + query_idx * 0.5),
                    Inches(9), Inches(0.4)
                )
                stats_frame = stats_box.text_frame
                stats_frame.text = placeholder
                stats_frame.paragraphs[0].font.size = Pt(14)
        
        os.makedirs(config.BASE_TEMPLATE_DIR, exist_ok=True)
        customer_dir = os.path.join(config.BASE_TEMPLATE_DIR, customer_name)
        os.makedirs(customer_dir, exist_ok=True)
        
        output_path = os.path.join(customer_dir, f"{customer_name}_auto_template.pptx")
        prs.save(output_path)
        
        results["generated_file"] = output_path
        results["summary"] = {
            "total_panels": len(panels),
            "selected_panels": len(selected_panels) if selected_panels else len(panels),
            "total_placeholders": len(results["placeholders"])
        }
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"템플릿 생성 실패: {str(e)}")
    
    return results
