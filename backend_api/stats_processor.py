import os
import re
import json
import time
import datetime
import requests
from dateutil.relativedelta import relativedelta
from pptx import Presentation
import config

def normalize_title(text):
    normalized = re.sub(r'[^a-zA-Z0-9가-힣]', '', text)
    return normalized.lower()

def find_all_templates(root_dir):
    templates = []
    if not os.path.isdir(root_dir):
        return templates
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.endswith('.pptx') and not filename.startswith('~$'):
                full_path = os.path.join(dirpath, filename)
                templates.append(full_path)
    return templates

def calculate_previous_month_dates():
    today = datetime.date.today()
    end_date = today.replace(day=1) - relativedelta(days=1)
    start_date = end_date.replace(day=1)
    start_date_str_kr = start_date.strftime("%Y년 %m월 %d일")
    end_date_str_kr = end_date.strftime("%Y년 %m월 %d일")
    start_date_str_hyphen = start_date.strftime("%Y-%m-%d")
    end_date_str_hyphen = end_date.strftime("%Y-%m-%d")
    return {
        "placeholders": {
            "{{START_DATE}}": start_date_str_kr,
            "{{END_DATE}}": end_date_str_kr,
            "{{MONTH}}": end_date.strftime("%m"),
            "{{DATE_RANGE}}": f"{start_date_str_kr} ~ {end_date_str_kr}",
            "{{DATE_RANGE_HYPHEN}}": f"{start_date_str_hyphen} ~ {end_date_str_hyphen}"
        },
        "filename_date": end_date.strftime("%Y-%m"),
        "start_ts": int(datetime.datetime.combine(start_date, datetime.time.min).timestamp() * 1000),
        "end_ts": int(datetime.datetime.combine(end_date, datetime.time.max).timestamp() * 1000)
    }

def find_all_panels_recursively(panel_list):
    all_panels = []
    for panel in panel_list:
        all_panels.append(panel)
        if panel.get("type") == "row" and "panels" in panel:
            all_panels.extend(find_all_panels_recursively(panel["panels"]))
    return all_panels

def get_dashboard_definition(dashboard_uid, retries=3, delay=1):
    url = f"{config.GRAFANA_URL.rstrip('/')}/api/dashboards/uid/{dashboard_uid}"
    headers = {"Authorization": f"Bearer {config.API_KEY}"}
    for i in range(retries):
        try:
            response = requests.get(url, headers=headers, timeout=20, verify=config.VERIFY_SSL)
            response.raise_for_status()
            dashboard = response.json().get('dashboard', {})
            dashboard['all_panels'] = find_all_panels_recursively(dashboard.get('panels', []))
            return dashboard
        except requests.exceptions.RequestException as e:
            if i < retries - 1:
                time.sleep(delay)
    return None

def find_panel_by_title(all_panels, title_from_placeholder):
    normalized_placeholder_title = normalize_title(title_from_placeholder)
    for panel in all_panels:
        panel_title_from_grafana = panel.get('title', '')
        normalized_grafana_title = normalize_title(panel_title_from_grafana)
        if normalized_grafana_title and normalized_grafana_title == normalized_placeholder_title:
            return panel
    return None

def get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts):
    headers = {"Authorization": f"Bearer {config.API_KEY}", "Content-Type": "application/json"}
    
    all_queries = panel.get('targets', [])
    selected = [q.copy() for q in all_queries if q.get('refId') == query_letter]
    if not selected:
        return None
    
    panel_datasource = panel.get('datasource')
    if panel_datasource:
        for q in selected:
            if 'datasource' not in q or q['datasource'] is None:
                q['datasource'] = panel_datasource
    
    for q in selected:
        q.pop('real_hosts', None)
    
    for q in selected:
        q.setdefault('maxDataPoints', 720)
        q.setdefault('intervalMs', 3600000)
    
    query_payload = {
        "queries": selected,
        "from": str(start_ts),
        "to": str(end_ts)
    }
    query_url = f"{config.GRAFANA_URL.rstrip('/')}/api/ds/query"
    
    try:
        response = requests.post(query_url, headers=headers,
                                 data=json.dumps(query_payload),
                                 timeout=120, verify=config.VERIFY_SSL)
        response.raise_for_status()
        return response.json()
    except requests.exceptions.RequestException:
        return None

def get_all_placeholders(prs):
    placeholders = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    matches = re.findall(r'(\{\{.*?\}\})', p.text)
                    for match in matches:
                        placeholders.add(match)
    return list(placeholders)

def replace_text_in_presentation(prs, replacements):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, value in replacements.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)

def process_statistics(customer_name=None):
    results = {
        "success": True,
        "processed_files": [],
        "errors": [],
        "failed_placeholders": [],
        "summary": {}
    }
    
    try:
        date_info = calculate_previous_month_dates()
        date_placeholders = date_info["placeholders"]
        filename_date = date_info["filename_date"]
        start_ts = date_info["start_ts"]
        end_ts = date_info["end_ts"]
        
        if customer_name:
            templates = []
            for root, dirs, files in os.walk(config.OUTPUT_DIR_WITH_IMAGES):
                rel_path = os.path.relpath(root, config.OUTPUT_DIR_WITH_IMAGES)
                if rel_path.startswith(customer_name) or rel_path == customer_name:
                    for f in files:
                        if f.endswith('.pptx') and not f.startswith('~$'):
                            templates.append(os.path.join(root, f))
            results["summary"]["mode"] = f"특정 고객사: {customer_name}"
        else:
            templates = find_all_templates(config.OUTPUT_DIR_WITH_IMAGES)
            results["summary"]["mode"] = "전체 고객사"
        
        if not templates:
            results["success"] = False
            results["errors"].append("처리할 템플릿이 없습니다.")
            return results
        
        results["summary"]["total_templates"] = len(templates)
        
        for template_path in templates:
            try:
                relative_path = os.path.relpath(os.path.dirname(template_path), config.OUTPUT_DIR_WITH_IMAGES)
                first_level_folder = relative_path.split(os.sep)[0] if relative_path != '.' else None
                
                dashboard_uid = None
                if first_level_folder and first_level_folder in config.DASHBOARD_MAP:
                    dashboard_uid = config.DASHBOARD_MAP[first_level_folder]
                
                if not dashboard_uid:
                    results["errors"].append(f"고객사 '{first_level_folder}' 대시보드 매핑 없음")
                    continue
                
                prs = Presentation(template_path)
                placeholders_in_template = get_all_placeholders(prs)
                
                grafana_placeholders = [ph for ph in placeholders_in_template 
                                       if re.match(r'\{\{[^{}]+_[A-Z]\}\}', ph)]
                
                if not grafana_placeholders:
                    replacements = date_placeholders.copy()
                    replace_text_in_presentation(prs, replacements)
                    
                    output_subdir = os.path.join(config.OUTPUT_DIR, relative_path)
                    os.makedirs(output_subdir, exist_ok=True)
                    final_output_path = os.path.join(output_subdir, os.path.basename(template_path))
                    prs.save(final_output_path)
                    
                    results["processed_files"].append({
                        "template": os.path.basename(template_path),
                        "customer": first_level_folder,
                        "output_path": final_output_path,
                        "grafana_queries": 0
                    })
                    continue
                
                dashboard = get_dashboard_definition(dashboard_uid)
                if not dashboard:
                    results["errors"].append(f"Grafana 대시보드 조회 실패: {dashboard_uid}")
                    continue
                
                all_panels = dashboard.get('all_panels', [])
                replacements = date_placeholders.copy()
                
                for placeholder in grafana_placeholders:
                    match = re.match(r'\{\{([^{}]+)_([A-Z])\}\}', placeholder)
                    if not match:
                        continue
                    
                    panel_title_from_placeholder = match.group(1)
                    query_letter = match.group(2)
                    
                    panel = find_panel_by_title(all_panels, panel_title_from_placeholder)
                    if not panel:
                        results["failed_placeholders"].append({
                            "placeholder": placeholder,
                            "reason": "패널을 찾을 수 없음"
                        })
                        continue
                    
                    stats_result = get_grafana_stats_by_panel(panel, query_letter, start_ts, end_ts)
                    if not stats_result:
                        results["failed_placeholders"].append({
                            "placeholder": placeholder,
                            "reason": "Grafana API 조회 실패"
                        })
                        continue
                    
                    values = []
                    for result_item in stats_result.get('results', {}).values():
                        for frame in result_item.get('frames', []):
                            for field in frame.get('schema', {}).get('fields', []):
                                if field.get('name') != 'Time':
                                    field_values = frame.get('data', {}).get('values', [[]])[frame.get('schema', {}).get('fields', []).index(field)]
                                    values.extend([v for v in field_values if v is not None])
                    
                    if values:
                        max_val = round(max(values), 2)
                        mean_val = round(sum(values) / len(values), 2)
                        replacements[placeholder] = config.SENTENCE_TEMPLATE.format(max=max_val, mean=mean_val)
                    else:
                        results["failed_placeholders"].append({
                            "placeholder": placeholder,
                            "reason": "데이터 없음"
                        })
                
                replace_text_in_presentation(prs, replacements)
                
                output_subdir = os.path.join(config.OUTPUT_DIR, relative_path)
                os.makedirs(output_subdir, exist_ok=True)
                final_output_path = os.path.join(output_subdir, os.path.basename(template_path))
                prs.save(final_output_path)
                
                results["processed_files"].append({
                    "template": os.path.basename(template_path),
                    "customer": first_level_folder,
                    "output_path": final_output_path,
                    "grafana_queries": len(grafana_placeholders)
                })
                
            except Exception as e:
                results["errors"].append(f"템플릿 처리 실패 ({os.path.basename(template_path)}): {str(e)}")
        
        results["summary"]["processed_count"] = len(results["processed_files"])
        results["summary"]["error_count"] = len(results["errors"])
        results["summary"]["failed_placeholder_count"] = len(results["failed_placeholders"])
        
        if results["summary"]["processed_count"] == 0:
            results["success"] = False
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"전체 프로세스 실패: {str(e)}")
    
    return results
