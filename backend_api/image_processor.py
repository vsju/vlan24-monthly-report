import os
import re
import datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches
import config

def normalize_name(name):
    normalized = re.sub(r'[^a-zA-Z0-9가-힣]', '', name)
    return normalized.lower()

def find_all_templates(root_dir):
    templates = []
    if not os.path.isdir(root_dir):
        return templates
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.endswith('.pptx') and not filename.startswith('~$'):
                full_path = os.path.join(dirpath, filename)
                templates.append(full_path)
    return templates

def find_templates_for_customer(template_base_dir, customer_path):
    customer_template_dir = os.path.join(template_base_dir, customer_path)
    return find_all_templates(customer_template_dir)

def calculate_previous_month_dates():
    today = datetime.date.today()
    end_date = today.replace(day=1) - relativedelta(days=1)
    start_date = end_date.replace(day=1)
    start_date_str_kr = start_date.strftime("%Y년 %m월 %d일")
    end_date_str_kr = end_date.strftime("%Y년 %m월 %d일")
    start_date_str_hyphen = start_date.strftime("%Y-%m-%d")
    end_date_str_hyphen = end_date.strftime("%Y-%m-%d")
    return {
        "placeholders": {
            "{{START_DATE}}": start_date_str_kr,
            "{{END_DATE}}": end_date_str_kr,
            "{{MONTH}}": end_date.strftime("%m"),
            "{{DATE_RANGE}}": f"{start_date_str_kr} ~ {end_date_str_kr}",
            "{{DATE_RANGE_HYPHEN}}": f"{start_date_str_hyphen} ~ {end_date_str_hyphen}"
        },
        "filename_date": end_date.strftime("%Y-%m")
    }

def find_available_images(root_dir):
    image_map = {}
    if not os.path.isdir(root_dir):
        return None
    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.gif')):
                name_without_extension = os.path.splitext(filename)[0]
                image_map[name_without_extension] = os.path.join(dirpath, filename)
    return image_map

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for placeholder, value in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)

def process_images(customer_name=None):
    results = {
        "success": True,
        "processed_files": [],
        "errors": [],
        "summary": {}
    }
    
    try:
        date_info = calculate_previous_month_dates()
        date_placeholders = date_info["placeholders"]
        filename_date = date_info["filename_date"]
        
        if customer_name:
            templates = find_templates_for_customer(config.BASE_TEMPLATE_DIR, customer_name)
            results["summary"]["mode"] = f"특정 고객사: {customer_name}"
        else:
            templates = find_all_templates(config.BASE_TEMPLATE_DIR)
            results["summary"]["mode"] = "전체 고객사"
        
        if not templates:
            results["success"] = False
            results["errors"].append("처리할 템플릿이 없습니다.")
            return results
        
        results["summary"]["total_templates"] = len(templates)
        
        for template_path in templates:
            try:
                prs = Presentation(template_path)
                
                relative_path = os.path.relpath(os.path.dirname(template_path), config.BASE_TEMPLATE_DIR)
                image_search_dir = os.path.join(config.BASE_IMAGE_DIR, relative_path)
                if not os.path.isdir(image_search_dir):
                    parent_dir = os.path.dirname(image_search_dir)
                    if os.path.isdir(parent_dir):
                        image_search_dir = parent_dir
                
                image_map = find_available_images(image_search_dir)
                
                if image_map is None:
                    results["errors"].append(f"이미지 폴더 없음: {image_search_dir}")
                    continue
                
                inserted_count = 0
                
                for slide in prs.slides:
                    for shape in slide.shapes:
                        replace_text_in_shape(shape, date_placeholders)
                        
                        shape_name_normalized = normalize_name(shape.name)
                        for img_name_original, img_path in image_map.items():
                            img_name_normalized = normalize_name(img_name_original)
                            if shape_name_normalized == img_name_normalized:
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height
                                
                                slide.shapes._spTree.remove(shape._element)
                                slide.shapes.add_picture(img_path, left, top, width, height)
                                inserted_count += 1
                                break
                
                output_subdir = os.path.join(config.OUTPUT_DIR_WITH_IMAGES, relative_path)
                if not os.path.exists(output_subdir):
                    os.makedirs(output_subdir)
                
                intermediate_output_path = os.path.join(
                    output_subdir,
                    os.path.basename(template_path)
                )
                prs.save(intermediate_output_path)
                
                results["processed_files"].append({
                    "template": os.path.basename(template_path),
                    "customer": relative_path if relative_path != '.' else 'root',
                    "images_inserted": inserted_count,
                    "output_path": intermediate_output_path
                })
                
            except Exception as e:
                results["errors"].append(f"템플릿 처리 실패 ({os.path.basename(template_path)}): {str(e)}")
        
        results["summary"]["processed_count"] = len(results["processed_files"])
        results["summary"]["error_count"] = len(results["errors"])
        
        if results["summary"]["processed_count"] == 0:
            results["success"] = False
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"전체 프로세스 실패: {str(e)}")
    
    return results
