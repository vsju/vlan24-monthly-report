import os
import requests
import json
import datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import config

def get_high_usage_metrics(dashboard_uid, threshold=90):
    url = f"{config.GRAFANA_URL.rstrip('/')}/api/dashboards/uid/{dashboard_uid}"
    headers = {"Authorization": f"Bearer {config.API_KEY}"}
    
    try:
        response = requests.get(url, headers=headers, timeout=20, verify=config.VERIFY_SSL)
        response.raise_for_status()
        dashboard = response.json().get('dashboard', {})
        
        today = datetime.date.today()
        end_date = today.replace(day=1) - relativedelta(days=1)
        start_date = end_date.replace(day=1)
        start_ts = int(datetime.datetime.combine(start_date, datetime.time.min).timestamp() * 1000)
        end_ts = int(datetime.datetime.combine(end_date, datetime.time.max).timestamp() * 1000)
        
        high_usage = []
        
        def extract_panels(panel_list):
            panels = []
            for panel in panel_list:
                if panel.get('type') != 'row':
                    panels.append(panel)
                if panel.get('type') == 'row' and 'panels' in panel:
                    panels.extend(extract_panels(panel['panels']))
            return panels
        
        panels = extract_panels(dashboard.get('panels', []))
        
        for panel in panels:
            panel_title = panel.get('title', '')
            panel_id = panel.get('id', '')
            
            if not any(keyword in panel_title.lower() for keyword in ['cpu', 'memory', 'disk', '메모리', '디스크']):
                continue
            
            targets = panel.get('targets', [])
            for target in targets:
                ref_id = target.get('refId', 'A')
                
                selected = [target.copy()]
                panel_datasource = panel.get('datasource')
                if panel_datasource:
                    for q in selected:
                        if 'datasource' not in q or q['datasource'] is None:
                            q['datasource'] = panel_datasource
                
                for q in selected:
                    q.pop('real_hosts', None)
                    q.setdefault('maxDataPoints', 720)
                    q.setdefault('intervalMs', 3600000)
                
                query_payload = {
                    "queries": selected,
                    "from": str(start_ts),
                    "to": str(end_ts)
                }
                query_url = f"{config.GRAFANA_URL.rstrip('/')}/api/ds/query"
                
                try:
                    query_response = requests.post(
                        query_url, 
                        headers={"Authorization": f"Bearer {config.API_KEY}", "Content-Type": "application/json"},
                        data=json.dumps(query_payload),
                        timeout=120, 
                        verify=config.VERIFY_SSL
                    )
                    query_response.raise_for_status()
                    stats_result = query_response.json()
                    
                    values = []
                    for result_item in stats_result.get('results', {}).values():
                        for frame in result_item.get('frames', []):
                            for field in frame.get('schema', {}).get('fields', []):
                                if field.get('name') != 'Time':
                                    field_values = frame.get('data', {}).get('values', [[]])[
                                        frame.get('schema', {}).get('fields', []).index(field)
                                    ]
                                    values.extend([v for v in field_values if v is not None])
                    
                    if values:
                        max_val = round(max(values), 2)
                        if max_val >= threshold:
                            high_usage.append({
                                "panel_id": panel_id,
                                "panel_title": panel_title,
                                "query": ref_id,
                                "max_usage": max_val,
                                "metric_type": "CPU" if "cpu" in panel_title.lower() else 
                                             "메모리" if any(k in panel_title.lower() for k in ['memory', '메모리']) else 
                                             "디스크" if any(k in panel_title.lower() for k in ['disk', '디스크']) else "기타"
                            })
                
                except Exception:
                    continue
        
        return high_usage
        
    except Exception as e:
        raise Exception(f"Grafana 메트릭 조회 실패: {str(e)}")

def add_recommendations(file_path, customer_name, threshold=90):
    results = {
        "success": True,
        "recommendations": [],
        "errors": []
    }
    
    try:
        if not os.path.exists(file_path):
            results["success"] = False
            results["errors"].append(f"파일을 찾을 수 없음: {file_path}")
            return results
        
        if customer_name not in config.DASHBOARD_MAP:
            results["success"] = False
            results["errors"].append(f"'{customer_name}' 대시보드 매핑이 없습니다.")
            return results
        
        dashboard_uid = config.DASHBOARD_MAP[customer_name]
        high_usage = get_high_usage_metrics(dashboard_uid, threshold)
        
        if not high_usage:
            results["recommendations"].append("모든 시스템이 정상 범위 내에서 운영되고 있습니다.")
        else:
            prs = Presentation(file_path)
            
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = "권고사항"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            
            content_top = 1.5
            for idx, item in enumerate(high_usage):
                recommendation_text = (
                    f"{idx + 1}. {item['panel_title']} (ID: {item['panel_id']})\n"
                    f"   - {item['metric_type']} 사용률: {item['max_usage']}%\n"
                    f"   - 권고: 사용률이 높은 상태입니다. 시스템 증설 또는 최적화를 고려하세요."
                )
                
                text_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(content_top),
                    Inches(9), Inches(1.2)
                )
                text_frame = text_box.text_frame
                text_frame.text = recommendation_text
                text_frame.paragraphs[0].font.size = Pt(14)
                
                results["recommendations"].append(recommendation_text)
                content_top += 1.3
                
                if content_top > 6.5:
                    slide = prs.slides.add_slide(blank_slide_layout)
                    content_top = 1.5
            
            prs.save(file_path)
        
        results["summary"] = {
            "high_usage_count": len(high_usage),
            "threshold": threshold
        }
        
    except Exception as e:
        results["success"] = False
        results["errors"].append(f"권고사항 생성 실패: {str(e)}")
    
    return results
